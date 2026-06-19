"""
PPT generator — clones slides from the reference template and fills placeholders.

Why clone instead of draw from scratch?
  The reference .pptx already has the brand fonts (Anton/Poppins), the colour
  palette, the canvas size, and the decorative graphics baked in.
  By cloning a template slide and only replacing the text we get a pixel-perfect
  match for free.

Theory slides are the one exception — the template has no theory layout, so we
clone the Recap layout (slide 1) and rewrite the big heading text with the
topic title. The numbered bullet boxes already match what a theory slide needs.

Canvas sizes vary across templates:
  • Common Template.pptx       — 40 × 22.5 in  (scale = 1.0)
  • CLAT / Architecture        — 10 × 5.6 in   (scale = 0.25)
All programmatic drawing coordinates are authored in the 40×22.5 reference space
and then multiplied by _style().scale before passing to pptx so every template
renders correctly without per-template branching.

Style extraction:
  Every template has a style-guide slide as its LAST slot.  parse_template_style()
  reads that slide automatically so any future template added to reference_ppts/
  is picked up without touching any Python code.
"""
import os
import re
import copy
import math
import threading
from dataclasses import dataclass, field
from pptx import Presentation
from pptx.util import Pt
from pptx.oxml.ns import qn
from lxml import etree

from schemas.slide_content import SlideContent
from schemas.slide_plan import TemplateType
from schemas.request import PDFContext
from config import OUTPUT_DIR, TEMPLATE_PPTX, DEVANAGARI_FONT
from pipeline.formula_renderer import (
    split_bullet_segments,
    is_pure_formula_bullet,
    render_formula_png,
    FormulaSegment,
    TextSegment,
)


# Devanagari block (U+0900–U+097F) — used to detect Hindi text runs that the
# brand Latin fonts (Anton/Poppins) cannot render.
_DEVANAGARI_RE = re.compile(r'[\u0900-\u097F]')


# ── Template style — read from every template's built-in style-guide slide ────
#
# Every reference PPTX has a "style guide" as its LAST slide (slot 13 in the
# 14-slot ordering).  That slide contains lines like:
#
#   Font Type  : Poppins Bold          ← heading / body font
#   Font Color : #FFCC31               ← brand accent (Cambria Math section)
#
# parse_template_style() reads these lines so NO Python changes are needed
# when a new template is dropped into reference_ppts/.
#
# Canvas dimensions are also read from the template so all programmatic
# drawing auto-scales:
#   Common Template  40 × 22.5 in  → scale 1.00
#   CLAT / Arch      10 × 5.60 in  → scale 0.25


@dataclass
class TemplateStyle:
    """Visual spec extracted from a template's style-guide slide (last slot)."""
    canvas_w:   float = 40.0     # slide width in inches
    canvas_h:   float = 22.5     # slide height in inches
    scale:      float = 1.0      # canvas_w / 40.0
    accent_hex: str   = "FFCC31" # brand accent (6 hex chars, no '#')
    body_font:  str   = "Poppins"


_STYLE_CACHE: dict[str, TemplateStyle] = {}


def parse_template_style(tpl_path: str) -> TemplateStyle:
    """
    Open *tpl_path*, read its last slide (style-guide), and return a
    TemplateStyle with canvas dimensions, accent colour, and body font.

    Results are cached by file path so the PPTX is only parsed once per process.
    Works automatically for any future template as long as it follows the
    convention of having its style-guide as the final slide.
    """
    if tpl_path in _STYLE_CACHE:
        return _STYLE_CACHE[tpl_path]

    try:
        prs = Presentation(tpl_path)
        canvas_w = prs.slide_width.inches
        canvas_h = prs.slide_height.inches
        scale    = canvas_w / 40.0

        # Collect all text from the style-guide slide (always the last one).
        style_slide = prs.slides[-1]
        all_text = " ".join(
            shape.text_frame.text
            for shape in style_slide.shapes
            if shape.has_text_frame
        )

        # Accent colour — style guides have MULTIPLE "Font Color" lines:
        #   • Normal text sections use "Font Color : #FFFFFF" (white)
        #   • The Cambria Math / equation section has "Font Color :#XXXX" —
        #     that's the brand accent colour (gold, amber, orange, etc.)
        # We look only at shapes whose runs use Cambria Math font.
        accent_hex = "FFCC31"
        for shape in style_slide.shapes:
            if not shape.has_text_frame:
                continue
            shape_text = shape.text_frame.text
            if "Font Color" not in shape_text:
                continue
            uses_cambria = any(
                run.font.name and "Cambria" in (run.font.name or "")
                for para in shape.text_frame.paragraphs
                for run in para.runs
            )
            if not uses_cambria:
                continue
            m = re.search(r'Font\s+Color\s*:\s*#([0-9A-Fa-f]{6})', shape_text)
            if m:
                accent_hex = m.group(1).upper()
                break

        # Body font: first non-Cambria "Font Type : <name>" entry.
        body_font = "Poppins"
        for m in re.finditer(r'Font\s+Type\s*:\s*([^\n|]+)', all_text):
            fname = m.group(1).strip()
            if "Cambria" not in fname and "Math" not in fname:
                body_font = fname.split()[0]   # e.g. "Poppins" from "Poppins Bold"
                break

        style = TemplateStyle(
            canvas_w=canvas_w, canvas_h=canvas_h, scale=scale,
            accent_hex=accent_hex, body_font=body_font,
        )
        print(
            f"  [style] {os.path.basename(tpl_path)}: "
            f"{canvas_w:.1f}×{canvas_h:.1f}in  scale={scale:.3f}  "
            f"accent=#{accent_hex}  font={body_font}"
        )
    except Exception as exc:
        import logging
        logging.getLogger(__name__).warning(
            "Could not parse template style from %s: %s", tpl_path, exc
        )
        style = TemplateStyle()

    _STYLE_CACHE[tpl_path] = style
    return style


# Thread-local variable so that concurrent asyncio.to_thread() generate calls
# each see their own TemplateStyle without interfering with each other.
_tl = threading.local()
_DEFAULT_STYLE = TemplateStyle()


def _style() -> TemplateStyle:
    """Return the TemplateStyle for the current generation thread."""
    return getattr(_tl, "style", _DEFAULT_STYLE)


def _accent_rgb():
    """Return pptx RGBColor for this thread's accent (default: yellow #FFCC31)."""
    from pptx.dml.color import RGBColor as _RGB
    h = _style().accent_hex
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    return _RGB(r, g, b)


def _insert_formula_image(
    slide,
    png_bytes: bytes,
    left_in: float,
    top_in: float,
    max_width_in: float = 30.0,
    max_height_in: float = 5.0,
) -> None:
    """
    Insert a rendered formula PNG onto the slide as a floating picture shape.

    The image is scaled to fit within (max_width_in × max_height_in) while
    preserving its aspect ratio. Transparent PNGs look great on the dark slide
    background.
    """
    import io as _io
    from pptx.util import Inches as _Inches

    # Measure natural size from PNG header via PIL
    try:
        from PIL import Image as _PILImage
        img = _PILImage.open(_io.BytesIO(png_bytes))
        nat_w_px, nat_h_px = img.size
    except Exception:
        nat_w_px, nat_h_px = 800, 200    # safe fallback

    # Scale to fit within bounding box, preserve aspect ratio
    aspect = nat_w_px / max(nat_h_px, 1)
    w_in = min(max_width_in, max_height_in * aspect)
    h_in = w_in / aspect
    if h_in > max_height_in:
        h_in = max_height_in
        w_in = h_in * aspect

    # Centre horizontally
    centre_in = left_in + max_width_in / 2
    img_left_in = centre_in - w_in / 2

    img_stream = _io.BytesIO(png_bytes)
    slide.shapes.add_picture(
        img_stream,
        _Inches(img_left_in), _Inches(top_in),
        width=_Inches(w_in), height=_Inches(h_in),
    )


def _force_run_font(run, name: str) -> None:
    """
    Force a run to use `name` for the Latin, East-Asian AND complex-script
    typeface slots.

    PowerPoint picks the COMPLEX-SCRIPT (a:cs) font for Devanagari, while
    LibreOffice may use a:latin — so we set all three to be safe. python-pptx's
    `run.font.name` only writes a:latin, so we add a:ea / a:cs at the XML level
    (they belong AFTER a:latin in the rPr schema order).
    """
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = etree.SubElement(rPr, qn(tag))
        el.set("typeface", name)


def _apply_devanagari_fonts(slide) -> None:
    """
    Final pass over a finished slide: any text run that contains Devanagari
    characters is re-assigned to the Devanagari font, since Anton/Poppins have
    no Hindi glyphs and would otherwise render as tofu boxes (□□□).

    Covers both normal text frames and table cells.
    """
    def _fix_text_frame(tf):
        for para in tf.paragraphs:
            for run in para.runs:
                if _DEVANAGARI_RE.search(run.text or ""):
                    _force_run_font(run, DEVANAGARI_FONT)

    for shape in slide.shapes:
        if shape.has_text_frame:
            _fix_text_frame(shape.text_frame)
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    _fix_text_frame(cell.text_frame)


# Decorative icons placed by layout. Each lives in backend/assets/visuals/.
_ASSETS_DIR   = os.path.dirname(os.path.dirname(TEMPLATE_PPTX))
_VISUALS_DIR  = os.path.join(_ASSETS_DIR, "visuals")
# recap.png is placed as a bottom-banner on programmatic summary slides
# (summary slides clone the blank dark, so they have no template illustration).
# The other visuals-folder icons (topic-heading, slide-heading, summary) are no
# longer added programmatically — template decorations come through via cloning.
_RECAP_ICON_PATH   = os.path.join(_VISUALS_DIR, "recap.png")


# ─────────────────────────────────────────────────────────────────────────────
# TEMPLATE INDEX MAP — which slide in Common Template.pptx is the source
# (0-based; template has 14 slides total)
# ─────────────────────────────────────────────────────────────────────────────
#  0  Recap of previous lecture (orange heading + 4 numbered points + decor)
#  1  Topics to be covered      (same layout, different decorative picture)
#  2  Section heading           ("Type Heading Here" — big centred text)
#  3  Blank content slide
#  4  MCQ — vertical 4 options (A/B/C/D stacked)
#  5  MCQ — 2x2 grid of options
#  6  Question only (no options)
#  7  PYQ MCQ — vertical options, wider "Question (PYQ Exam-Year)" bar
#  8  PYQ MCQ — 2x2 grid
#  9  PYQ Question only
# 10  Summary  ("Summary" small heading + decor)
# 11  Homework ("Homework" small heading + decor)
# 12  Thank You (BLANK layout — decorative)
# 13  Style guide reference (skip)

LAYOUT_TO_TEMPLATE_IDX = {
    TemplateType.title_slide:        2,   # use section heading style for title
    TemplateType.recap_slide:        0,
    TemplateType.topics_slide:       1,
    TemplateType.section_heading:    2,
    TemplateType.theory_slide:       3,   # blank dark slide — we draw heading + bullets
    TemplateType.table_slide:        3,   # blank dark slide — we draw heading + table
    TemplateType.theory_table_slide: 3,   # blank dark slide — heading + bullets + table
    TemplateType.passage_slide:      3,   # blank dark slide — we draw banner + passage
    TemplateType.mcq_slide:          4,
    TemplateType.mcq_grid_slide:     5,
    TemplateType.question_only:      6,
    TemplateType.pyq_slide:          7,
    TemplateType.pyq_grid_slide:     8,
    TemplateType.pyq_question_only:  9,
    TemplateType.summary:           10,
    TemplateType.homework_slide:    11,
    TemplateType.thank_you_slide:   12,
    TemplateType.figure_slide:       3,   # blank dark slide — we draw tag + image/text
}


# ─────────────────────────────────────────────────────────────────────────────
# XML helpers — clone & delete slides at the OOXML level
# ─────────────────────────────────────────────────────────────────────────────

def _clone_slide(prs, src_slide):
    """
    Deep-copy a slide's shape tree AND background image into a new slide.

    The template's slide backgrounds are embedded images (dark theme +
    decorative panels) referenced via relationship IDs (rId*). If we only
    cloned shape XML, the new slide would inherit the slide master's
    default WHITE background and our white body text would become invisible.

    So we:
      1. Copy every shape from the source slide.
      2. Copy the source slide's <p:bg> element.
      3. Walk the source slide's relationships and copy any referenced
         media parts (images) into the new slide's relationships, fixing
         up rId references inside the cloned XML.
    """
    blank_layout = prs.slide_layouts[0]
    new_slide = prs.slides.add_slide(blank_layout)

    new_cSld    = new_slide._element.find(qn('p:cSld'))
    src_cSld    = src_slide._element.find(qn('p:cSld'))
    new_sptree  = new_cSld.find(qn('p:spTree'))
    src_sptree  = src_cSld.find(qn('p:spTree'))

    # ── 1. Reset new slide's shape tree and copy shapes ─────────────────────
    for child in list(new_sptree):
        if etree.QName(child).localname not in ('nvGrpSpPr', 'grpSpPr'):
            new_sptree.remove(child)
    for child in src_sptree:
        if etree.QName(child).localname in ('nvGrpSpPr', 'grpSpPr'):
            continue
        new_sptree.append(copy.deepcopy(child))

    # ── 2. Copy <p:bg> from source slide (background image / fill) ──────────
    src_bg = src_cSld.find(qn('p:bg'))
    if src_bg is not None:
        # Remove any existing bg on the new slide first
        existing_bg = new_cSld.find(qn('p:bg'))
        if existing_bg is not None:
            new_cSld.remove(existing_bg)
        # <p:bg> must come BEFORE <p:spTree> per OOXML schema
        new_cSld.insert(list(new_cSld).index(new_sptree), copy.deepcopy(src_bg))

    # ── 3. Migrate referenced parts (background image + any embedded media) ─
    # Find every r:embed / r:link attribute inside the cloned XML and rebind
    # those rIds to fresh relationships on the new slide part.
    R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    src_part = src_slide.part
    new_part = new_slide.part
    rid_map: dict[str, str] = {}

    def _ensure_rid(old_rid: str) -> str | None:
        if old_rid in rid_map:
            return rid_map[old_rid]
        try:
            rel = src_part.rels[old_rid]
        except KeyError:
            return None
        # Relate the new slide part to the SAME target part as the source
        new_rid = new_part.relate_to(rel.target_part, rel.reltype)
        rid_map[old_rid] = new_rid
        return new_rid

    for elem in new_cSld.iter():
        for attr_name in (f"{{{R_NS}}}embed", f"{{{R_NS}}}link"):
            old_rid = elem.get(attr_name)
            if old_rid:
                new_rid = _ensure_rid(old_rid)
                if new_rid:
                    elem.set(attr_name, new_rid)

    return new_slide


def _delete_slides_by_indices(prs, indices):
    """Drop slides at the given 0-based indices (delete in reverse order)."""
    xml_slides = prs.slides._sldIdLst
    slides_list = list(xml_slides)
    for idx in sorted(indices, reverse=True):
        if 0 <= idx < len(slides_list):
            sl_el = slides_list[idx]
            rid = sl_el.get(qn('r:id'))
            xml_slides.remove(sl_el)
            prs.part.drop_rel(rid)


# ─────────────────────────────────────────────────────────────────────────────
# Text-replacement helpers
# ─────────────────────────────────────────────────────────────────────────────

def _iter_runs(slide):
    """Yield every text run on the slide."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                yield run


def _replace_first(slide, placeholder, new_value):
    """Replace the FIRST occurrence of `placeholder` with `new_value`."""
    if not new_value:
        return False
    for run in _iter_runs(slide):
        if placeholder in run.text:
            run.text = run.text.replace(placeholder, new_value)
            return True
    return False


def _replace_sequence(slide, placeholder, values):
    """
    Replace each occurrence of `placeholder` with the next item from `values`.
    Used for 4-option MCQ / numbered bullets — each is its own textbox.
    """
    it = iter(values)
    for run in _iter_runs(slide):
        if placeholder in run.text:
            try:
                run.text = run.text.replace(placeholder, next(it))
            except StopIteration:
                # leave remaining placeholders blank to avoid stray "Type option here"
                run.text = run.text.replace(placeholder, "")


def _replace_placeholders_by_shape_position(slide, placeholder, values, key):
    """
    Fill placeholder textboxes in a SPECIFIC visual order, not the XML shape order.
    `key` is a callable receiving each shape and returning a sort key.
    Used for grid layouts where XML shape order != visual order (A, C, B, D).
    """
    targets = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if placeholder in shape.text_frame.text:
            targets.append(shape)

    targets.sort(key=key)
    for shape, value in zip(targets, values):
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if placeholder in run.text:
                    run.text = run.text.replace(placeholder, value)

    # blank any remaining unfilled placeholders
    for shape in targets[len(values):]:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                run.text = run.text.replace(placeholder, "")


def _grid_position_key(shape):
    """Sort key giving visual A, B, C, D order on a 2x2 grid."""
    # row-major: top-row first (smaller top), then left to right
    return (round(shape.top, -5), round(shape.left, -5))


def _clear_unused_placeholders(slide):
    """Blank out any leftover 'Type ... here' text so it doesn't show in output."""
    for run in _iter_runs(slide):
        if "Type option here" in run.text or "Type question here" in run.text:
            run.text = ""
        if "Type Heading Here" in run.text:
            run.text = ""


def _resolve_font_pt(shape, fallback_pt: int) -> int:
    """Return the first explicit font size on a shape, or fallback."""
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.font.size:
                return int(run.font.size.pt)
    return fallback_pt


def _apply_heading_style(
    shape,
    text_len: int,
    base_pt: int,
    min_pt: int,
    color,
    wrap: bool = True,
    max_width_in: float | None = None,
    char_width_factor: float = 0.00568,
):
    """
    Apply a heading font size that keeps the rendered text within bounds.

    When `max_width_in` is given, compute a font size such that
        text_len × pt × char_width_factor ≤ max_width_in
    so the heading text doesn't overflow its visual budget. The default
    `char_width_factor` (0.00568 in/char/pt) is calibrated against Anton at
    264pt for "Recap" (5 chars ≈ 7.5 in). Use ~0.0080 for wider fonts like
    Poppins. The template defaults (264pt big / 132pt sub) only fit the
    original "Recap" / "of previous lecture" text — any other content must
    be width-fit, not just length-banded.
    """
    from pptx.util import Pt

    if not shape:
        return
    target_pt = base_pt
    if max_width_in and text_len > 0:
        fit_pt = int(max_width_in / (text_len * char_width_factor))
        target_pt = max(min(base_pt, fit_pt), min_pt)
    elif text_len > 80:
        target_pt = max(base_pt - 28, min_pt)
    elif text_len > 65:
        target_pt = max(base_pt - 16, min_pt)
    elif text_len > 50:
        target_pt = max(base_pt - 8, min_pt)

    shape.text_frame.word_wrap = wrap
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(target_pt)
            run.font.color.rgb = color


# ─────────────────────────────────────────────────────────────────────────────
# Explicit logo cleanup
# ─────────────────────────────────────────────────────────────────────────────

# pw_badge_top_right.png and the old _add_top_right_badge / _remove_explicit_top_left_logo
# functions have been removed.  Every reference template already has the PW badge
# embedded in its slides; the clone operation brings it across automatically.
# Adding it again from a PNG file produced duplicate badges (the duplicate-check
# tolerance was smaller than the positional offset between template badge and the
# programmatic position, so the check silently failed).


# ─────────────────────────────────────────────────────────────────────────────
# Per-template fillers
# ─────────────────────────────────────────────────────────────────────────────

def _strip_question_prefix(text: str) -> str:
    """Drop 'Q:', 'Question:', 'Q.1', '1.' etc. from the front of a question."""
    t = text.strip()
    for prefix in ("Question:", "Question.", "Q:", "Q.", "Ques:", "Ques.", "Problem:", "Q "):
        if t.startswith(prefix):
            t = t[len(prefix):].strip()
    # also drop leading numbering like "1." or "1) "
    if t and t[0].isdigit():
        i = 0
        while i < len(t) and (t[i].isdigit() or t[i] in ".)"):
            i += 1
        t = t[i:].strip()
    return t


def _strip_option_prefix(text: str) -> str:
    """Drop '(a)', 'a)', 'A.', etc. from each option."""
    t = text.strip()
    for pfx in (
        "(a) ", "(b) ", "(c) ", "(d) ",
        "(A) ", "(B) ", "(C) ", "(D) ",
        "a) ", "b) ", "c) ", "d) ",
        "A) ", "B) ", "C) ", "D) ",
        "a. ", "b. ", "c. ", "d. ",
        "A. ", "B. ", "C. ", "D. ",
    ):
        if t.startswith(pfx):
            return t[len(pfx):].strip()
    return t


def _fill_recap_or_topics(slide, content: SlideContent):
    """
    Slides 1/2 of template share the same shape pattern:
      - "Recap" / "Topics"      (big Anton, orange)        — first run
      - "of previous lecture" / "to be covered" (Arial)    — second run
      - 4× "Type option here"                              — bullet textboxes
    We override the heading with the slide title (split into two visual lines
    for nicer fit), and feed key_points into the bullets.
    """
    title = content.title.strip()
    # Render the whole title as ONE uniform-size heading line (all words the
    # same big size, naturally spaced) instead of a big first word + small
    # remainder. `small` is cleared so the sub-heading box stays empty.
    big = title
    small = ""

    # Find the two heading textboxes — they're the first two shapes that contain
    # the original placeholder words "Recap"/"Topics" and "of previous"/"to be".
    from pptx.util import Pt
    from pptx.dml.color import RGBColor

    heading_set = False
    sub_set = False
    heading_shape = None
    sub_shape = None
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        full = shape.text_frame.text.strip()
        if full in ("Recap", "Topics") and not heading_set:
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    r.text = big
            heading_set = True
            heading_shape = shape
        elif full.startswith(("of previous", "to be")) and not sub_set:
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    r.text = small
            sub_set = True
            sub_shape = shape

    YELLOW = _accent_rgb()  # template accent colour
    base_big = _resolve_font_pt(heading_shape, 90) if heading_shape else 90
    base_small = _resolve_font_pt(sub_shape, 40) if sub_shape else 40

    # Note: no programmatic icon is added here.  Recap/Topics slides (indexes 0/1)
    # already carry all decorative artwork from the reference template — the clone
    # brings it across automatically.  Adding extra icons from the visuals folder
    # would double-up the decoration and break template-specific styling.

    # The whole title now lives in the heading box on one line. Give it a wide
    # budget up to a safe right limit that stays clear of the right-side artwork,
    # and widen the box so it never wraps.
    from pptx.util import Inches
    RIGHT_LIMIT = 20.0
    heading_left_in = (heading_shape.left or 0) / 914400.0 if heading_shape else 1.63
    heading_max_w = max(RIGHT_LIMIT - heading_left_in, 6.0)
    if heading_shape is not None:
        heading_shape.width = Inches(heading_max_w + 1.0)

    # Single line, no wrap: shrink the font so the whole title fits the width.
    # char_width_factor ~0.0080 matches the heading font's real glyph width.
    _apply_heading_style(
        heading_shape,
        text_len=len(big),
        base_pt=base_big,
        min_pt=44,
        color=YELLOW,
        wrap=False,
        max_width_in=heading_max_w,
        char_width_factor=0.0080,
    )
    _apply_heading_style(
        sub_shape,
        text_len=len(small),
        base_pt=base_small,
        min_pt=36,
        color=YELLOW,
        wrap=True,
        max_width_in=12.0,
        char_width_factor=0.0080,
    )

    # ── THEORY SLIDE DECOR CLEANUP ───────────────────────────────────────────
    # Theory slides reuse the Recap layout (src_idx 0). That layout has a 
    # large decorative picture (books/calculator) at the top-right.
    # We remove it for theory slides to keep them clean and professional.
    if content.layout == TemplateType.theory_slide:
        # The decorative picture in Slide 0 sits at approx (23.43, 2.37)
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        PIC_LEFT = int(23.43 * 914400)
        PIC_TOP  = int(2.37 * 914400)
        TOL      = int(1.0 * 914400)
        for shape in list(slide.shapes):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                if abs(shape.left - PIC_LEFT) < TOL and abs(shape.top - PIC_TOP) < TOL:
                    # Remove the shape from the slide's shape tree
                    sp = shape._element
                    sp.getparent().remove(sp)

    _replace_sequence(slide, "Type option here", content.bullets[:4])
    _clear_unused_placeholders(slide)


def _find_backdrop_pill(slide, text_shape):
    """
    Find the decorative rounded-rect (the teal/blue "pill") that sits BEHIND the
    heading text box.

    On the section-heading template the blue pill is a SEPARATE auto-shape — the
    text box itself is transparent. To widen the blue background for a long
    heading we must resize that shape, so we locate it as the smallest auto-shape
    whose bounding box contains the text box's centre and is at least as tall as
    the text box (this ignores the small circular logo badge to its left).
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    if text_shape.left is None or text_shape.width is None:
        return None

    cx = text_shape.left + text_shape.width / 2
    cy = text_shape.top + text_shape.height / 2

    best = None
    best_area = None
    for sh in slide.shapes:
        if sh is text_shape:
            continue
        if sh.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
            continue
        if sh.left is None or sh.top is None or sh.width is None or sh.height is None:
            continue
        within = (sh.left <= cx <= sh.left + sh.width) and (sh.top <= cy <= sh.top + sh.height)
        if not within or sh.height < text_shape.height:
            continue
        area = sh.width * sh.height
        if best is None or area < best_area:
            best = sh
            best_area = area
    return best


def _fill_section_heading(slide, content: SlideContent):
    """
    Replace 'Type Heading Here' with the section title, keeping it on ONE LINE
    that always sits inside the blue rounded-rect pill.

    The blue pill is a SEPARATE decorative auto-shape behind a transparent text
    box. Wrapping the title to a 2nd/3rd line is what used to push text out of
    the pill — so we never wrap. Instead, for longer titles we:

      1. Leave short titles exactly as the template designed them.
      2. EXPAND the blue pill (and its text box) to the right — up to a bound
         that clears the top-right badge — so the bigger heading fits on one
         line inside the blue background.
      3. Only if the widest the pill can grow STILL isn't enough, SHRINK the
         font so the whole title fits that single line.

    This is applied per-heading, so only the long ones grow / shrink.
    """
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    text = (content.title or "").strip()
    _replace_first(slide, "Type Heading Here", text)

    YELLOW = _accent_rgb()  # template accent colour

    target_shape = None
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if shape.text_frame.text.strip() == text:
            target_shape = shape
            break

    if target_shape is None or not text:
        _clear_unused_placeholders(slide)
        return

    base_pt = _resolve_font_pt(target_shape, 84)
    text_len = max(len(text), 1)

    sc = _style().scale
    # Conservative glyph width — LibreOffice substitutes a wider bold sans for
    # the brand font on export, so text renders longer than a naive estimate.
    CHAR_W = 0.0115
    EDGE_PAD_IN = 1.1 * sc              # clearance from the pill's rounded ends
    RIGHT_BOUND_IN = 35.4 * sc          # keep clear of the top-right PW badge
    _ABS_MIN_TITLE_PT = max(int(36 * sc), 8)   # last-resort readable floor

    tb_top_in = target_shape.top / 914400
    tb_h_in   = target_shape.height / 914400

    pill = _find_backdrop_pill(slide, target_shape)
    if pill is not None:
        pill_left_in = pill.left / 914400
        pill_w_in    = pill.width / 914400
        pill_top_in  = pill.top / 914400
        pill_h_in    = pill.height / 914400
    else:
        pill_left_in = target_shape.left / 914400
        pill_w_in    = target_shape.width / 914400
        pill_top_in  = tb_top_in
        pill_h_in    = tb_h_in

    # Width the title needs on ONE line at the template's base font.
    needed_w = text_len * base_pt * CHAR_W + 2 * EDGE_PAD_IN
    max_pill_w = max(RIGHT_BOUND_IN - pill_left_in, pill_w_in)

    chosen_pt = base_pt
    if needed_w <= pill_w_in:
        # Short title — already fits the existing pill; leave geometry as-is.
        new_pill_w = pill_w_in
    elif needed_w <= max_pill_w:
        # Longer title — EXPAND the blue pill to fit it at full size.
        new_pill_w = needed_w
    else:
        # Too long even for the widest pill — expand to the max, then SHRINK
        # the font so the whole title still fits on this one line.
        new_pill_w = max_pill_w
        usable = max(new_pill_w - 2 * EDGE_PAD_IN, 1.0)
        chosen_pt = max(int(usable / (text_len * CHAR_W)), _ABS_MIN_TITLE_PT)

    # Apply the new width to the blue pill (keep its left/top), then lay the text
    # box over it with padding so the single line centres inside the blue area.
    if pill is not None and abs(new_pill_w - pill_w_in) > 0.01:
        pill.width = Inches(new_pill_w)

    target_shape.left = Inches(pill_left_in + EDGE_PAD_IN)
    target_shape.width = Inches(max(new_pill_w - 2 * EDGE_PAD_IN, 1.0))
    target_shape.top = Inches(pill_top_in)
    target_shape.height = Inches(pill_h_in)

    tf = target_shape.text_frame
    tf.word_wrap = False          # never wrap — the title stays on one line
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for para in tf.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        for run in para.runs:
            run.font.size = Pt(chosen_pt)
            run.font.color.rgb = YELLOW

    # Note: no programmatic icon is added here.  The section-heading slide (index 2)
    # already has decorative icons baked into every reference template; they come
    # through automatically when the slide is cloned, so adding another one from
    # the visuals folder would double-up the decoration inconsistently.
    _clear_unused_placeholders(slide)


_SUBBULLET_RE = None  # lazy-built in _fill_theory_slide


def _strip_theory_prefix(text: str) -> str:
    """Drop the writer-injected '-> ' (or '➤ ') marker so the renderer owns the arrow."""
    t = text.strip()
    for pfx in ("-> ", "->", "➤ ", "➤", "• ", "•"):
        if t.startswith(pfx):
            return t[len(pfx):].lstrip()
    return t


def _clear_bullet_props(pPr):
    """Remove any existing bullet/indent child elements so we can re-set them."""
    from pptx.oxml.ns import qn
    for tag in ("a:buClr", "a:buClrTx", "a:buSzPct", "a:buSzPts", "a:buSzTx",
                "a:buFont", "a:buFontTx", "a:buChar", "a:buNone", "a:buAutoNum",
                "a:tabLst"):
        for el in pPr.findall(qn(tag)):
            pPr.remove(el)


def _set_arrow_bullet(paragraph, indent_in: float, color_hex: str = "FFCC31"):
    """
    Apply a NATIVE PowerPoint bullet (➤) with a hanging indent.

    Native bullets are the reliable way to get the behaviour a hand-made deck
    has: the arrow sits at the left, the text starts at `indent_in`, and every
    WRAPPED line aligns with the text — not under the arrow. LibreOffice honours
    this where a manual "arrow run + tab" does not.

    marL = indent_in (text + wrapped lines), indent = -indent_in (bullet hangs).
    """
    from pptx.util import Inches
    from pptx.oxml.ns import qn

    pPr = paragraph._p.get_or_add_pPr()
    marL = int(Inches(indent_in))
    pPr.set("marL", str(marL))
    pPr.set("indent", str(-marL))
    _clear_bullet_props(pPr)
    # order matters (schema): buClr, buFont, buChar — appended after spcAft.
    bu_clr = etree.SubElement(pPr, qn("a:buClr"))
    etree.SubElement(bu_clr, qn("a:srgbClr")).set("val", color_hex)
    etree.SubElement(pPr, qn("a:buFont")).set("typeface", "Arial")
    etree.SubElement(pPr, qn("a:buChar")).set("char", "➤")


def _set_plain_hanging(paragraph, indent_in: float, hang_in: float):
    """Hanging indent with NO bullet glyph — for sub-bullets whose '(a)' is text."""
    from pptx.util import Inches
    from pptx.oxml.ns import qn

    pPr = paragraph._p.get_or_add_pPr()
    pPr.set("marL", str(int(Inches(indent_in))))
    pPr.set("indent", str(-int(Inches(hang_in))))
    _clear_bullet_props(pPr)
    etree.SubElement(pPr, qn("a:buNone"))


def _is_solution_slide(content: SlideContent) -> bool:
    """Detect if this theory slide is a solution/worked-example slide."""
    t = (content.title or "").strip().lower()
    return t.startswith("solution") and (len(t) < 9 or t[8:9] in (':', ' ', '-', '—'))


def _fill_theory_slide(slide, content: SlideContent, strategy=None):
    """
    Theory / concept layout built on top of the blank dark slide (template idx 3).

    Layout:
      - Compact rounded-rect tag at top-left holding the title in bold.
        Yellow for normal theory; green (#4CAF50) for solution slides.
        Width auto-fits the title length.
      - Body textbox with arrow (➤) bullets in white. Bullets prefixed with
        "(a) ", "(b) ", "(c) ", "(d) " auto-indent as sub-bullets without arrow.

    The writer's "-> " prefix is stripped here so the visual arrow stays the
    renderer's responsibility.
    """
    import re
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE

    is_solution = _is_solution_slide(content)

    YELLOW = _accent_rgb()  # template accent colour
    GREEN  = RGBColor(0x4C, 0xAF, 0x50)
    TAG_BG = GREEN if is_solution else YELLOW
    BLACK  = RGBColor(0x10, 0x10, 0x10)
    WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
    TAG_TEXT = WHITE if is_solution else BLACK

    sub_re = re.compile(r'^\(\s*([a-dA-D])\s*\)\s*')

    # Canvas scale — 1.0 for Common Template (40×22.5 in), 0.25 for CLAT/Arch (10×5.6 in).
    # Every Inches() / Pt() call below is already expressed in the 40×22.5 reference
    # space; multiplying by `sc` maps them onto whatever canvas the chosen template uses.
    sc = _style().scale

    # ── Title — colored tag, auto-sized to text ───────────────────────────────
    raw_title = (content.title or "").strip()
    if is_solution:
        display = raw_title
        for pfx in ("Solution:", "Solution -", "Solution —", "Solution "):
            if display.lower().startswith(pfx.lower()):
                display = "SOLUTION: " + display[len(pfx):].strip().upper()
                break
        else:
            display = raw_title.upper()
        title = display
    else:
        title = raw_title.upper() if raw_title else "TOPIC"

    # _fit_title_tag() returns values already scaled to the current canvas.
    PAD_X = 0.6 * sc
    PAD_Y = 0.20 * sc
    title_pt, tag_w, tag_h = _fit_title_tag(title)   # already scaled
    tag_l = 1.0 * sc
    tag_t = 0.8 * sc
    # Note: no programmatic icon is added here.  The template's blank-dark slide
    # (index 3) has no decorative elements by design; any branding icons on other
    # slide types come through automatically when that slide is cloned.

    tag = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(tag_l), Inches(tag_t), Inches(tag_w), Inches(tag_h),
    )
    tag.fill.solid()
    tag.fill.fore_color.rgb = TAG_BG
    tag.line.fill.background()
    tag.shadow.inherit = False
    tag.adjustments[0] = 0.12

    tf = tag.text_frame
    tf.word_wrap = False
    try:
        from pptx.enum.text import MSO_AUTO_SIZE
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    except Exception:
        pass
    tf.margin_left = Inches(PAD_X)
    tf.margin_right = Inches(PAD_X)
    tf.margin_top = Inches(PAD_Y)
    tf.margin_bottom = Inches(PAD_Y)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.size = Pt(title_pt)    # title_pt already scaled by _fit_title_tag
    run.font.bold = True
    run.font.name = "Anton"
    run.font.color.rgb = TAG_TEXT

    # ── Body — arrow bullets + optional (a)/(b) sub-bullets ──────────────────
    body_left_in = 1.5 * sc
    body_top_in  = tag_t + tag_h + 0.9 * sc          # all already scaled
    body_width_in = 37.0 * sc
    canvas_h = _style().canvas_h                       # actual canvas height (in)
    body_height_in = canvas_h - body_top_in - 1.2 * sc   # leave footer room

    body_left   = Inches(body_left_in)
    body_top    = Inches(body_top_in)
    body_width  = Inches(body_width_in)
    body_height = Inches(max(body_height_in, 0.5 * sc))

    bullets = [b for b in (content.bullets or []) if b and b.strip()]
    # Drop bullets that are EMPTY once the "-> "/"➤" marker is removed — these
    # would otherwise render as a lone arrow with no text.
    bullets = [b for b in bullets if _strip_theory_prefix(b).strip()]
    if not bullets:
        return

    # Body font size: the fit engine returns a size calibrated for the 40×22.5
    # reference canvas; scale it down proportionally for smaller templates.
    from pipeline.fit_engine import consistent_body_font
    body_pt = max(int(consistent_body_font(TemplateType.theory_slide, strategy) * sc), 6)

    # ── Pre-render formula bullets ────────────────────────────────────────────
    # Bullets that are pure $$latex$$ blocks are rendered to PNG and inserted
    # as picture shapes instead of text runs. We track their position index so
    # we can compute the approximate Y coordinate after the text-frame loop.
    formula_positions: list[tuple[int, bytes]] = []   # (bullet_index, png_bytes)
    for bi, raw in enumerate(bullets):
        if is_pure_formula_bullet(raw):
            segs = split_bullet_segments(raw)
            for seg in segs:
                if isinstance(seg, FormulaSegment):
                    png = render_formula_png(
                        seg.latex,
                        fontsize=int(body_pt * 0.85),
                        text_color="white",
                        bg_transparent=True,
                    )
                    if png:
                        formula_positions.append((bi, png))
                    break   # only one formula per bullet line

    body_tb = slide.shapes.add_textbox(body_left, body_top, body_width, body_height)
    bt = body_tb.text_frame
    bt.word_wrap = True
    bt.vertical_anchor = MSO_ANCHOR.TOP

    # Per-bullet height estimate (scaled inches) for Y-offset computation.
    # body_pt is already scaled, so (body_pt / 72) gives scaled inches directly.
    space_main_pt = max(int(24 * sc), 4)
    space_sub_pt  = max(int(14 * sc), 3)
    sub_pt        = max(body_pt - max(int(4 * sc), 1), 6)
    _MAIN_BULLET_H = (body_pt * 1.25 + space_main_pt) / 72.0   # scaled inches
    _SUB_BULLET_H  = (sub_pt  * 1.25 + space_sub_pt)  / 72.0
    _FORMULA_SLOT_H = 4.0 * sc    # reserved slot for a formula image (scaled in)

    # Cumulative Y tracker (scaled inches, relative to body_top_in)
    cumulative_y: list[float] = []  # Y start of each bullet
    y_cursor = 0.0

    first = True
    formula_indices = {bi for bi, _ in formula_positions}

    for bi, raw in enumerate(bullets):
        cumulative_y.append(y_cursor)

        if bi in formula_indices:
            # Skip: formula bullet — rendered as image below, reserve slot
            y_cursor += _FORMULA_SLOT_H
            if first:
                first = False
            continue

        text = _strip_theory_prefix(raw)
        is_sub = bool(sub_re.match(text))

        p = bt.paragraphs[0] if first else bt.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(space_sub_pt if is_sub else space_main_pt)

        # Hanging-indent: proportional to font size so arrow→text gap is
        # readable at any canvas scale.
        main_indent = round(body_pt * 0.020, 3)   # naturally scaled via body_pt

        if is_sub:
            _set_plain_hanging(p, main_indent + 0.7 * sc, 0.55 * sc)
            run_t = p.add_run()
            run_t.text = text
            run_t.font.size = Pt(sub_pt)
            run_t.font.name = "Poppins"
            run_t.font.color.rgb = WHITE
            y_cursor += _SUB_BULLET_H
        else:
            _set_arrow_bullet(p, main_indent, color_hex=_style().accent_hex)
            run_t = p.add_run()
            run_t.text = text
            run_t.font.size = Pt(body_pt)
            run_t.font.name = "Poppins"
            run_t.font.color.rgb = YELLOW
            y_cursor += _MAIN_BULLET_H

    # ── Insert formula images at their estimated Y positions ─────────────────
    for bi, png_bytes in formula_positions:
        slot_top = body_top_in + cumulative_y[bi]
        _insert_formula_image(
            slide,
            png_bytes,
            left_in=body_left_in,
            top_in=slot_top,
            max_width_in=36.0 * sc,
            max_height_in=max(_FORMULA_SLOT_H - 0.4 * sc, 0.2),
        )


# ─────────────────────────────────────────────────────────────────────────────
# Table renderers — table_slide and theory_table_slide
# ─────────────────────────────────────────────────────────────────────────────

def _fit_title_tag(title: str) -> tuple[int, float, float]:
    """
    Pick (font_pt, tag_width_in, tag_height_in) so the title NEVER overflows
    its rounded-rect background.

    All returned values are **already scaled** to the current template's canvas
    (via _style().scale), so callers pass them directly to Pt() / Inches() with
    no further multiplication.

    Reference sizes (scale = 1.0, i.e. Common Template 40×22.5 in):
      PAD_X = 0.6 in, PAD_Y = 0.20 in, USABLE_MAX_W = 33 in,
      title_pt candidates: 54 → 48 → 42 → 36 → 32 → 28 pt (fixed baseline).

    CONSISTENCY: every title starts from the same TARGET size so the vast
    majority of headings — which are short — all render at exactly TARGET.
    We only step the font DOWN for a genuinely long title.
    """
    sc = _style().scale           # e.g. 1.0 for Common, 0.25 for CLAT/Arch
    PAD_X = 0.6 * sc
    PAD_Y = 0.20 * sc
    SAFETY = 0.5 * sc
    USABLE_MAX_W = 33.0 * sc      # hard right bound before the PW badge
    char_w = 0.0135               # conservative for the bold substitute font
    budget = USABLE_MAX_W - 2 * PAD_X - SAFETY

    n = max(len(title), 1)
    # Fixed target so headings are consistent; only shrink when too long to fit.
    title_pt = max(int(28 * sc), 6)
    for candidate_pt in (54, 48, 42, 36, 32, 28):
        scaled_pt = max(int(candidate_pt * sc), 6)
        if n * scaled_pt * char_w <= budget:
            title_pt = scaled_pt
            break

    text_w_in = n * title_pt * char_w
    tag_w = max(
        min(text_w_in + 2 * PAD_X + SAFETY, USABLE_MAX_W),
        max(5.0 * sc, 1.0),
    )
    tag_h = title_pt / 72.0 + 2 * PAD_Y + 0.35 * sc
    return title_pt, tag_w, tag_h


def _draw_yellow_title_tag(slide, raw_title: str, top_in: float = 0.8,
                           is_solution: bool = False) -> tuple[float, float]:
    """
    Shared helper: draws a rounded-rect title tag at the top of the slide.
    Yellow for normal slides, green for solution slides.

    `top_in` should be passed in the **40×22.5 reference space** (unscaled).
    Returns (tag_left_in, bottom_in) in the **scaled** canvas space so callers
    can position the body below the tag correctly.
    """
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE

    sc = _style().scale

    YELLOW = _accent_rgb()
    GREEN  = RGBColor(0x4C, 0xAF, 0x50)
    BLACK  = RGBColor(0x10, 0x10, 0x10)
    WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
    TAG_BG = GREEN if is_solution else YELLOW
    TAG_TEXT = WHITE if is_solution else BLACK

    title = (raw_title or "TOPIC").strip().upper() or "TOPIC"

    PAD_X = 0.6 * sc
    PAD_Y = 0.20 * sc
    title_pt, tag_w, tag_h = _fit_title_tag(title)   # already scaled
    tag_l = 1.0 * sc
    top_scaled = top_in * sc

    tag = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(tag_l), Inches(top_scaled), Inches(tag_w), Inches(tag_h),
    )
    tag.fill.solid()
    tag.fill.fore_color.rgb = TAG_BG
    tag.line.fill.background()
    tag.shadow.inherit = False
    tag.adjustments[0] = 0.12

    tf = tag.text_frame
    tf.word_wrap = False
    try:
        from pptx.enum.text import MSO_AUTO_SIZE
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    except Exception:
        pass
    tf.margin_left = Inches(PAD_X)
    tf.margin_right = Inches(PAD_X)
    tf.margin_top = Inches(PAD_Y)
    tf.margin_bottom = Inches(PAD_Y)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.size = Pt(title_pt)
    run.font.bold = True
    run.font.name = "Anton"
    run.font.color.rgb = TAG_TEXT

    return tag_l, top_scaled + tag_h


def _looks_numeric(value: str) -> bool:
    """Heuristic: 'is this cell a number?' — used for default column alignment."""
    if value is None:
        return False
    s = str(value).strip()
    if not s:
        return False
    # strip common decorations: currency, %, parentheses, commas, leading +/-
    s = s.replace(",", "").replace("$", "").replace("₹", "").replace("%", "")
    s = s.replace("(", "").replace(")", "").strip()
    if not s:
        return False
    try:
        float(s)
        return True
    except ValueError:
        return False


def _pick_table_font_size(rows_count: int, cols_count: int,
                          available_height_in: float,
                          longest_cell_len: int,
                          available_col_width_in: float) -> int:
    """
    Pick a font size that lets `rows_count` rows fit in `available_height_in`
    AND lets the longest cell fit within `available_col_width_in`.

    Returns a pt size in [12, 32].
    """
    # Height-based cap: each row is roughly font_pt * 1.6 (incl. inner padding)
    # in points → inches = pt / 72.
    if rows_count <= 0:
        return 22
    h_pt = (available_height_in * 72.0) / (rows_count * 1.6)

    # Width-based cap: longest cell text shouldn't overflow its column.
    char_w_factor = 0.0080  # Poppins approx, inches per (char × pt)
    if longest_cell_len > 0 and available_col_width_in > 0:
        w_pt = available_col_width_in / (longest_cell_len * char_w_factor)
    else:
        w_pt = 999

    candidate = int(min(h_pt, w_pt))
    return max(12, min(candidate, 32))


def _add_styled_table(
    slide,
    headers: list[str],
    rows: list[list[str]],
    left_in: float,
    top_in: float,
    width_in: float,
    height_in: float,
    column_alignments: list[str] | None = None,
    vcenter: bool = True,
):
    """
    Render a real PowerPoint table on the dark slide.

    Visual language:
      - Header row : solid yellow (#FFCC31) fill, black bold text.
      - Body rows  : alternating very-dark fills (#1F1F1F / #2A2A2A) so each
                     row reads cleanly against the dark template background.
                     White text.
      - Borders    : thin dark grey lines so the grid is visible but not loud.
      - Column widths: content-proportional but CLAMPED so one long label
                       column can't dominate and starve the data columns.
      - Row heights : sized NATURALLY from the font (≈ comfortable padding),
                      not stretched to fill the slide. A small table stays a
                      small, tidy block instead of ballooning to full height.
      - Placement   : the finished block is VERTICALLY CENTERED within the
                      [top_in, top_in + height_in] region (when `vcenter`),
                      exactly like a hand-made deck would place it.

    `height_in` is the AVAILABLE region height (a budget), not a forced size.
    `column_alignments` is an optional list with "left"/"center"/"right" per
    column. When omitted, numeric-looking columns default to right-aligned
    and text columns default to left-aligned.
    """
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    YELLOW   = RGBColor(0xFF, 0xCC, 0x31)
    BLACK    = RGBColor(0x10, 0x10, 0x10)
    WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
    ROW_DARK = RGBColor(0x1F, 0x1F, 0x1F)
    ROW_MID  = RGBColor(0x2A, 0x2A, 0x2A)
    BORDER   = RGBColor(0x55, 0x55, 0x55)

    cols = len(headers)
    if cols == 0 or not rows:
        return
    # Defensive: pad/trim every row to len(headers).
    norm_rows = []
    for r in rows:
        r = list(r) + [""] * max(0, cols - len(r))
        norm_rows.append([str(c) if c is not None else "" for c in r[:cols]])

    total_rows = len(norm_rows) + 1  # +1 for header

    # ── Column widths — content-proportional but CLAMPED ───────────────────
    # A long row-label ("P.V.A.F for 4 year") shouldn't be 3-4× the numeric
    # columns. Clamp each column's effective length to [avg×0.7, avg×1.7] so
    # the table reads balanced; long labels simply word-wrap.
    col_text_len = [max(len(headers[c]), 1) for c in range(cols)]
    for r in norm_rows:
        for c in range(cols):
            col_text_len[c] = max(col_text_len[c], len(r[c]))
    avg_len = sum(col_text_len) / cols
    clamped = [min(max(l, avg_len * 0.7), avg_len * 1.7) for l in col_text_len]
    total_clamped = sum(clamped) or 1
    widths = [width_in * (l / total_clamped) for l in clamped]

    # ── Font size — fit the longest cell to its column width ───────────────
    longest_cell = max(col_text_len)
    avg_col_w = sum(widths) / cols
    font_pt = _pick_table_font_size(
        rows_count=total_rows,
        cols_count=cols,
        available_height_in=height_in,
        longest_cell_len=longest_cell,
        available_col_width_in=avg_col_w,
    )

    # ── Natural row heights (NOT stretched to fill) ────────────────────────
    # A comfortable row ≈ font height × 1.9 + small padding. Header a touch
    # taller. Cap so a few-row table never balloons; scale down only if the
    # natural block would exceed the available region.
    MIN_BODY_ROW_H = 0.65
    MAX_BODY_ROW_H = 1.35
    body_row_h = min(MAX_BODY_ROW_H, max(MIN_BODY_ROW_H, font_pt / 72.0 * 1.9 + 0.16))
    header_row_h = body_row_h * 1.2
    natural_h = header_row_h + body_row_h * (total_rows - 1)

    if natural_h > height_in:
        shrink = height_in / natural_h
        body_row_h *= shrink
        header_row_h *= shrink
        natural_h = height_in

    # Position the block within the available region. We bias toward the TOP
    # (capped gap) instead of dead-centering, so a short table sits just below
    # the title/caption rather than floating in the middle of a large void.
    place_top = top_in
    if vcenter and natural_h < height_in:
        gap = (height_in - natural_h) / 2.0
        place_top = top_in + min(gap, 1.6)

    table_shape = slide.shapes.add_table(
        total_rows, cols,
        Inches(left_in), Inches(place_top),
        Inches(width_in), Inches(natural_h),
    )
    table = table_shape.table

    for c in range(cols):
        table.columns[c].width = Inches(widths[c])

    table.rows[0].height = Inches(header_row_h)
    for r in range(1, total_rows):
        table.rows[r].height = Inches(body_row_h)

    # ── Default alignment per column ---------------------------------------
    if column_alignments and len(column_alignments) == cols:
        col_align = column_alignments
    else:
        col_align = []
        for c in range(cols):
            # numeric if >= 60% of body cells in this column parse as numbers
            n = sum(1 for r in norm_rows if _looks_numeric(r[c]))
            col_align.append("right" if (n / max(len(norm_rows), 1)) >= 0.6 else "left")

    def _apply_align(p, name):
        if name == "right":
            p.alignment = PP_ALIGN.RIGHT
        elif name == "center":
            p.alignment = PP_ALIGN.CENTER
        else:
            p.alignment = PP_ALIGN.LEFT

    # ── Header row ---------------------------------------------------------
    for c in range(cols):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = YELLOW
        cell.margin_left = Inches(0.10)
        cell.margin_right = Inches(0.10)
        cell.margin_top = Inches(0.06)
        cell.margin_bottom = Inches(0.06)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        tf = cell.text_frame
        tf.word_wrap = True
        tf.clear()
        p = tf.paragraphs[0]
        _apply_align(p, "center")  # headers always centered
        run = p.add_run()
        run.text = headers[c] or ""
        run.font.size = Pt(font_pt)
        run.font.bold = True
        run.font.name = "Poppins"
        run.font.color.rgb = BLACK

    # ── Body rows ----------------------------------------------------------
    for r_idx, row in enumerate(norm_rows):
        bg = ROW_DARK if r_idx % 2 == 0 else ROW_MID
        for c in range(cols):
            cell = table.cell(r_idx + 1, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            cell.margin_left = Inches(0.10)
            cell.margin_right = Inches(0.10)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            tf = cell.text_frame
            tf.word_wrap = True
            tf.clear()
            p = tf.paragraphs[0]
            _apply_align(p, col_align[c])
            run = p.add_run()
            run.text = row[c] or ""
            run.font.size = Pt(max(font_pt - 1, 12))
            run.font.name = "Poppins"
            run.font.color.rgb = WHITE

    return table_shape


def _draw_table_caption(slide, caption: str, left_in: float, top_in: float,
                        width_in: float) -> float:
    """Optional small italic caption above the table. Returns its bottom Y."""
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    if not caption:
        return top_in
    tb = slide.shapes.add_textbox(
        Inches(left_in), Inches(top_in), Inches(width_in), Inches(0.6),
    )
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = caption.strip()
    run.font.size = Pt(24)
    run.font.italic = True
    run.font.name = "Poppins"
    run.font.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    return top_in + 0.55


def _fill_table_slide(slide, content: SlideContent):
    """
    Table-only slide on the blank dark template (idx 3).

    Layout:
      - Yellow rounded-rect title tag at top-left (same style as theory_slide)
      - Optional italic caption immediately below the tag
      - A real PowerPoint table filling the rest of the body area

    Falls back to `_fill_theory_slide` (using the slide's bullets) if the
    writer didn't supply structured table_data — that way the slide is never
    empty even when the writer fails.
    """
    table_data = getattr(content, "table_data", None)
    if not table_data or not table_data.headers or not table_data.rows:
        # No usable table — let theory slide handle it as a graceful degrade.
        _fill_theory_slide(slide, content)
        return

    _, body_top = _draw_yellow_title_tag(
        slide, content.title, top_in=0.8,
        is_solution=_is_solution_slide(content),
    )

    LEFT = 1.0
    WIDTH = 38.0
    BOTTOM_LIMIT = 21.4

    cur_top = body_top + 0.45
    cur_top = _draw_table_caption(
        slide, table_data.caption or "", LEFT, cur_top, WIDTH,
    )

    table_h = max(BOTTOM_LIMIT - cur_top, 3.0)
    _add_styled_table(
        slide,
        headers=list(table_data.headers),
        rows=[list(r) for r in table_data.rows],
        left_in=LEFT,
        top_in=cur_top,
        width_in=WIDTH,
        height_in=table_h,
        column_alignments=table_data.column_alignments,
    )


def _fill_theory_table_slide(slide, content: SlideContent, strategy=None):
    """
    Theory bullets ABOVE a small table on the blank dark template (idx 3).

    Layout (top → bottom):
      1. Yellow title tag
      2. Theory bullets (arrow ➤ style, same as theory_slide)
      3. Optional table caption
      4. Real PowerPoint table

    The renderer guarantees no overlap: bullet block height is bounded so the
    table always has at least ~6 inches of vertical room. If the bullets would
    push the table off the slide, the body font shrinks first; if the table is
    still too tall, the renderer downgrades to a table-only slide (drops the
    bullets) so the data — which is harder to compress — stays readable.
    """
    import re
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    WHITE = RGBColor(0xFF, 0xFF, 0xFF)

    table_data = getattr(content, "table_data", None)
    bullets = [b for b in (content.bullets or []) if b and b.strip()]

    # If there's no table at all, fall through to the plain theory layout.
    if not table_data or not table_data.headers or not table_data.rows:
        _fill_theory_slide(slide, content, strategy)
        return

    # If there are no bullets, prefer the cleaner table-only layout.
    if not bullets:
        _fill_table_slide(slide, content)
        return

    _, body_top = _draw_yellow_title_tag(
        slide, content.title, top_in=0.8,
        is_solution=_is_solution_slide(content),
    )

    sc = _style().scale
    BODY_LEFT    = 1.5 * sc
    BODY_WIDTH   = 37.0 * sc
    BOTTOM_LIMIT = 21.4 * sc

    # Bullets block — give the bullets a bounded chunk of vertical space so
    # they never crowd the table. With at most 3 bullets at 32pt (after fit),
    # ~3.5in is plenty.
    bullets_top   = body_top + 0.5 * sc
    BULLETS_MAX_H = 5.0 * sc
    MIN_TABLE_H   = 6.0 * sc

    available_after_bullets = BOTTOM_LIMIT - (bullets_top + BULLETS_MAX_H + 0.6 * sc)
    if available_after_bullets < MIN_TABLE_H:
        BULLETS_MAX_H = max(2.5 * sc, BOTTOM_LIMIT - bullets_top - MIN_TABLE_H - 0.6 * sc)

    sub_re = re.compile(r'^\(\s*([a-dA-D])\s*\)\s*')

    bullet_pt = max(int((28 if len(bullets) >= 3 else 32) * sc), 6)

    body_tb = slide.shapes.add_textbox(
        Inches(BODY_LEFT), Inches(bullets_top),
        Inches(BODY_WIDTH), Inches(BULLETS_MAX_H),
    )
    bt = body_tb.text_frame
    bt.word_wrap = True
    bt.vertical_anchor = MSO_ANCHOR.TOP

    first = True
    for raw in bullets[:3]:
        text = _strip_theory_prefix(raw)
        is_sub = bool(sub_re.match(text))
        p = bt.paragraphs[0] if first else bt.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(max(int((10 if is_sub else 18) * sc), 3))
        main_indent = round(bullet_pt * 0.020, 3)
        if is_sub:
            _set_plain_hanging(p, main_indent + 0.7 * sc, 0.55 * sc)
            run_t = p.add_run()
            run_t.text = text
            run_t.font.size = Pt(max(bullet_pt - max(int(4 * sc), 1), 6))
            run_t.font.name = "Poppins"
            run_t.font.color.rgb = WHITE
        else:
            _set_arrow_bullet(p, main_indent, color_hex=_style().accent_hex)
            run_t = p.add_run()
            run_t.text = text
            run_t.font.size = Pt(bullet_pt)
            run_t.font.name = "Poppins"
            run_t.font.color.rgb = YELLOW

    # Table sits below the bullets block.
    table_top = bullets_top + BULLETS_MAX_H + 0.4 * sc
    table_top = _draw_table_caption(
        slide, table_data.caption or "", BODY_LEFT, table_top, BODY_WIDTH,
    )
    table_h = max(BOTTOM_LIMIT - table_top, MIN_TABLE_H)

    _add_styled_table(
        slide,
        headers=list(table_data.headers),
        rows=[list(r) for r in table_data.rows],
        left_in=BODY_LEFT,
        top_in=table_top,
        width_in=BODY_WIDTH,
        height_in=table_h,
        column_alignments=table_data.column_alignments,
    )


def _fill_passage_slide(slide, content: SlideContent):
    """
    Cloze / reading-comprehension passage on the blank dark slide (template idx 3).

    Visual language matches the theory slide (same color theme):
      - A yellow rounded-rect banner at top holding the `directions` line in
        black bold (e.g. "Directions (Q. 22-24): Cloze Test – Passage 1").
      - The verbatim passage paragraph below in white, justified, with every
        blank (__X__, .....(1).....) preserved EXACTLY as written by the writer.

    The passage is rendered as ONE flowing, atomic paragraph — never split or
    paraphrased — so the gaps stay visible for the student to fill.
    """
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from pipeline.fit_engine import estimate_block_height_in

    YELLOW = _accent_rgb()  # template accent colour
    BLACK  = RGBColor(0x10, 0x10, 0x10)
    WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

    # ── Directions banner — yellow rounded rect, black bold, word-wrapped ─────
    directions = (content.directions or content.title or "Passage").strip()

    band_l, band_t = 1.0, 0.8
    band_w = 36.0                       # wide banner; leaves room for PW badge
    band_pad_x, band_pad_y = 0.5, 0.18

    # Fit the directions font to the banner width (word-wrap to 1-2 lines).
    char_w = 0.0095                     # in/char/pt for the bold sans banner
    usable_w = band_w - 2 * band_pad_x
    dir_pt = 40
    for candidate_pt in (40, 36, 32, 28, 24):
        chars_per_line = max(1, int(usable_w / (candidate_pt * char_w)))
        if len(directions) <= chars_per_line * 2:   # fits within two lines
            dir_pt = candidate_pt
            break
    else:
        dir_pt = 24
    chars_per_line = max(1, int(usable_w / (dir_pt * char_w)))
    dir_lines = max(1, math.ceil(len(directions) / chars_per_line))
    band_h = (dir_pt / 72.0) * dir_lines * 1.25 + 2 * band_pad_y

    band = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(band_l), Inches(band_t), Inches(band_w), Inches(band_h),
    )
    band.fill.solid()
    band.fill.fore_color.rgb = YELLOW
    band.line.fill.background()
    band.shadow.inherit = False
    band.adjustments[0] = 0.18

    btf = band.text_frame
    btf.word_wrap = True
    btf.margin_left = Inches(band_pad_x)
    btf.margin_right = Inches(band_pad_x)
    btf.margin_top = Inches(band_pad_y)
    btf.margin_bottom = Inches(band_pad_y)
    btf.vertical_anchor = MSO_ANCHOR.MIDDLE
    bp = btf.paragraphs[0]
    bp.alignment = PP_ALIGN.LEFT
    brun = bp.add_run()
    brun.text = directions
    brun.font.size = Pt(dir_pt)
    brun.font.bold = True
    brun.font.name = "Anton"
    brun.font.color.rgb = BLACK

    # ── Passage body — white, justified, verbatim (blanks preserved) ──────────
    passage = (content.passage_text or "").strip()
    if not passage:
        # Fallback: writer didn't populate passage_text — join any bullets so
        # nothing is lost (still verbatim, just not pre-formatted).
        passage = "\n".join(b for b in (content.bullets or []) if b and b.strip())
    if not passage:
        return

    sc = _style().scale
    body_width_in  = 37.0 * sc
    body_top_in    = (band_t + band_h + 0.7 * sc)
    canvas_h       = _style().canvas_h
    body_height_in = max(6.0 * sc, canvas_h - body_top_in - 1.2 * sc)

    body_left = Inches(1.5 * sc)
    body_top  = Inches(body_top_in)

    # FILL THE SLIDE like a hand-made deck — not small text clustered at the top:
    #   1) pick the LARGEST font that still fits (reference range 30–56pt, scaled);
    #   2) spread the lines (line spacing) to use the leftover height;
    #   3) vertically CENTRE so any residual gap is balanced, never all at bottom.
    paras = [ln for ln in passage.split("\n") if ln.strip()] or [passage]
    # font range in reference space (40in); scale to current canvas
    pass_pt = max(int(30 * sc), 6)
    for pt_ref in range(56, 29, -2):
        pt = max(int(pt_ref * sc), 6)
        if estimate_block_height_in(paras, pt, body_width_in) <= body_height_in:
            pass_pt = pt
            break

    # Spread lines to fill the height.
    natural = estimate_block_height_in(paras, pass_pt, body_width_in)
    line_spacing = 1.3
    if natural > 0:
        line_spacing = max(1.3, min(1.55, 1.3 * (body_height_in / natural)))

    body_tb = slide.shapes.add_textbox(
        body_left, body_top, Inches(body_width_in), Inches(body_height_in)
    )
    bt = body_tb.text_frame
    bt.word_wrap = True
    bt.vertical_anchor = MSO_ANCHOR.MIDDLE

    for i, para_text in enumerate(paras):
        p = bt.paragraphs[0] if i == 0 else bt.add_paragraph()
        p.alignment = PP_ALIGN.JUSTIFY
        p.line_spacing = line_spacing
        p.space_after = Pt(max(int(18 * sc), 3))
        run = p.add_run()
        run.text = para_text
        run.font.size = Pt(pass_pt)
        run.font.name = "Poppins"
        run.font.color.rgb = WHITE


def _add_context_footer(slide, context: PDFContext):
    """
    Small context strip bottom-left on body slides — keeps subject / batch /
    purpose visible across the deck without crowding the template's layout.
    """
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    parts = [p for p in (context.subject, context.batch, context.purpose) if p]
    if not parts:
        return
    text = "  ·  ".join(parts)

    tb = slide.shapes.add_textbox(
        Inches(0.5), Inches(21.6), Inches(30.0), Inches(0.7)
    )
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(20)
    run.font.italic = True
    run.font.name = "Poppins"
    run.font.color.rgb = RGBColor(0x99, 0x99, 0x99)


def _fill_title_slide(slide, content: SlideContent, context: PDFContext):
    """
    Build a proper title slide on top of the section-heading layout.
    The template only has one heading textbox, so we add three more textboxes
    around it to show the full lecture context:
      - main heading  : the lecture / topic title  (replaces 'Type Heading Here')
      - subtitle      : Subject  ·  Batch
      - metadata      : Purpose  ·  Class Level  ·  Language
    """
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    YELLOW = _accent_rgb()  # template accent colour
    WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
    GRAY   = RGBColor(0xAA, 0xAA, 0xAA)

    # 1. Main heading — the lecture title
    topic = content.title or f"{context.subject} — {context.purpose}"
    _replace_first(slide, "Type Heading Here", topic)
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if shape.text_frame.text.strip() == topic:
            base_pt = _resolve_font_pt(shape, 84)
            _apply_heading_style(
                shape,
                text_len=len(topic),
                base_pt=base_pt,
                min_pt=56,
                color=YELLOW,
                wrap=True,
            )
            break

    # 2. Subtitle — Subject · Batch
    subtitle_parts = []
    if context.subject:
        subtitle_parts.append(context.subject)
    if context.batch:
        subtitle_parts.append(context.batch)
    subtitle = "  ·  ".join(subtitle_parts)

    if subtitle:
        tb = slide.shapes.add_textbox(
            Inches(1.0), Inches(10.0), Inches(38.0), Inches(2.0)
        )
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = subtitle
        run.font.size = Pt(60)
        run.font.bold = True
        run.font.name = "Poppins"
        run.font.color.rgb = WHITE

    # 3. Metadata — Purpose · Class Level · Language
    meta_parts = []
    if context.purpose:
        meta_parts.append(context.purpose)
    if context.class_level:
        meta_parts.append(context.class_level)
    if context.language:
        meta_parts.append(context.language)
    metadata = "    ·    ".join(meta_parts)

    if metadata:
        tb = slide.shapes.add_textbox(
            Inches(1.0), Inches(13.5), Inches(38.0), Inches(1.2)
        )
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = metadata
        run.font.size = Pt(36)
        run.font.name = "Poppins"
        run.font.color.rgb = GRAY

    # 4. Bottom accent bar — yellow line
    from pptx.util import Emu
    bar = slide.shapes.add_shape(
        1,  # rectangle
        Inches(8.0), Inches(20.5), Inches(24.0), Inches(0.15)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = YELLOW
    bar.line.fill.background()

    _clear_unused_placeholders(slide)


def _sanitize_question_title(text: str) -> str:
    """
    Safety net: strip any 'Answer:', 'Exam:', exam year tags, or explanation
    text that the writer may have accidentally put in the title field.
    """
    import re
    t = text.strip()
    # Cut at "Answer:" if present — everything after is the answer
    t = re.split(r'\s*Answer\s*:', t, maxsplit=1, flags=re.IGNORECASE)[0].strip()
    # Cut at "Exam:" if present — that's exam metadata
    t = re.split(r'\s*Exam\s*:', t, maxsplit=1, flags=re.IGNORECASE)[0].strip()
    # Cut at standalone "(SSC" or "(Exam" pattern — year tag in parens
    t = re.split(r'\s*\(\s*(?:SSC|Exam|PYQ|JEE|NEET|UPSC)', t, maxsplit=1)[0].strip()
    return t


def _shapes_with_text(slide, needle: str):
    """Return shapes whose text frame contains `needle`."""
    out = []
    for shape in slide.shapes:
        if shape.has_text_frame and needle in shape.text_frame.text:
            out.append(shape)
    return out


def _set_textbox_font_size(shape, size_pt: int) -> None:
    """Apply one font size to every run in a text box."""
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(size_pt)


def _estimated_wrapped_lines(text: str, font_pt: int, width_in: float) -> int:
    """
    Conservative line estimate for Poppins on the 40in canvas.
    Used only to decide whether a question needs extra vertical room.
    """
    chars_per_line = max(18, int(width_in / (font_pt * 0.0095)))
    return max(1, math.ceil(len(text) / chars_per_line))


# Easy knobs for long vertical MCQ question text. If a future long
# Assertion/Reason question looks too small/big, change these two values.
LONG_MCQ_QUESTION_FONT_PT = 60
VERY_LONG_MCQ_QUESTION_FONT_PT = 60
LONG_MCQ_QUESTION_THRESHOLD = 170
VERY_LONG_MCQ_QUESTION_THRESHOLD = 260
LONG_MCQ_MAX_QUESTION_HEIGHT_IN = 7.0

# ── Easy knobs for LONG vertical MCQ OPTIONS ─────────────────────────────────
# These only kick in when at least one option wraps to a 3rd line at the
# template's default option width. Short options keep the exact template look.
#   • WIDTH:   default template option box = 15.82in. We widen it a *little* so
#              a 3-line option becomes 2 lines. Increase for fewer lines.
#   • TRIGGER: how many estimated lines (at default width) before we adapt.
#   • GAP:     vertical breathing space between options when we reflow them.
#   • FONT:    options stay at template 54pt; only shrinks as a last resort if
#              the reflowed options would run off the bottom of the slide.
MCQ_OPTION_DEFAULT_WIDTH_IN = 15.82   # template width (don't change)
MCQ_OPTION_WIDENED_WIDTH_IN = 24.0    # ← widen long options to this (try 20–26)
MCQ_OPTION_TRIGGER_LINES    = 3       # ← act only when an option hits this many lines
MCQ_OPTION_FONT_PT          = 54      # template option font
MCQ_OPTION_GAP_IN           = 0.55    # ← extra gap between options when reflowing
MCQ_OPTION_LINE_HEIGHT_FACTOR = 1.2   # line height multiple (leave as is)
MCQ_OPTION_BAND_BOTTOM_IN   = 21.2    # keep options above this line
MCQ_OPTION_MIN_FONT_PT      = 44      # last-resort shrink floor

# ── Long-stem MCQ → 2-column option grid (slide-19 type) ─────────────────────
# When a *vertical* MCQ has a stem so long that 4 stacked options would collide
# with it, AND all four options are short enough to live in a half-width column,
# lay the options out as a 2x2 grid just below the stem:
#       A   B
#       C   D
# This halves the vertical room the options need, so everything stays on ONE
# slide. It is deliberately conservative — it only fires for genuinely long
# stems (>= MCQ_GRID_MIN_QUESTION_LINES wrapped lines). Normal vertical MCQs,
# A/R questions, grids and PYQ slides are untouched.
MCQ_GRID_QUESTION_THRESHOLD   = LONG_MCQ_QUESTION_THRESHOLD  # min stem chars (170)
MCQ_GRID_MIN_QUESTION_LINES   = 5     # only when the stem wraps to >= this many lines
MCQ_GRID_OPTION_MAX_LINES     = 2     # every option must fit <= this many lines (half width)
MCQ_GRID_MAX_QUESTION_HEIGHT_IN = 11.0  # let the long stem be this tall (uncapped vs vertical)
MCQ_GRID_TOP_GAP_IN           = 0.5   # gap between stem bottom and first option row
MCQ_GRID_ROW_GAP_IN           = 0.6   # gap between the two option rows
MCQ_GRID_COL_GUTTER_IN        = 0.6   # gap between the left/right columns
MCQ_GRID_BAND_BOTTOM_IN       = 21.2  # keep the bottom option row above this line


def _apply_long_vertical_mcq_layout(slide, question_shape, option_shapes, question: str) -> None:
    """
    Long Assertion/Reason MCQs can wrap into the fixed option area. For those
    rare cases only, shrink the question slightly, give it more height, and
    push the whole option band (labels + option text) downward.
    """
    if not question_shape or not option_shapes:
        return

    from pptx.util import Inches
    from pptx.enum.text import MSO_ANCHOR

    EMU_PER_IN = 914400
    q_len = len(question)
    if q_len < LONG_MCQ_QUESTION_THRESHOLD:
        return

    q_width_in = question_shape.width / EMU_PER_IN

    # A/R questions like the screenshot need smaller text and a taller question
    # area. Normal MCQs never enter this branch, so the template look is kept.
    q_font = (
        LONG_MCQ_QUESTION_FONT_PT
        if q_len < VERY_LONG_MCQ_QUESTION_THRESHOLD
        else VERY_LONG_MCQ_QUESTION_FONT_PT
    )
    lines = _estimated_wrapped_lines(question, q_font, q_width_in)
    needed_q_h = min(
        max(lines * q_font * 1.22 / 72.0 + 0.35, 2.6),
        LONG_MCQ_MAX_QUESTION_HEIGHT_IN,
    )

    question_shape.height = Inches(needed_q_h)
    question_shape.text_frame.word_wrap = True
    question_shape.text_frame.vertical_anchor = MSO_ANCHOR.TOP
    _set_textbox_font_size(question_shape, q_font)

    first_option_top = min(s.top for s in option_shapes)
    desired_option_top = question_shape.top + question_shape.height + Inches(0.35)
    shift = max(0, desired_option_top - first_option_top)
    if shift <= 0:
        return

    # Move the whole option band, not only the text placeholders. This keeps the
    # teal A/B/C/D badges aligned with their option text.
    band_top = first_option_top - Inches(0.80)
    band_bottom = Inches(19.8)
    max_shift = max(0, band_bottom - max(s.top + s.height for s in option_shapes))
    shift = min(shift, max_shift)
    if shift <= 0:
        return

    for shape in slide.shapes:
        if shape.top is None or shape.height is None:
            continue
        if band_top <= shape.top <= band_bottom:
            # Keep footer/logo/top question area untouched; option labels and
            # option text all live in this vertical band.
            shape.top = shape.top + shift


def _apply_long_vertical_mcq_options_layout(slide, option_shapes, options) -> None:
    """
    For vertical MCQs whose options are long (one or more options wrap to a 3rd
    line at the template width), widen the option boxes a little and re-flow them
    vertically so each option gets the room it needs without colliding with the
    next one. Short options are left exactly as the template has them.

    The A/B/C/D teal badges live as separate shapes to the LEFT of each option
    text box, so when an option moves we move its badge by the same amount to
    keep them aligned.
    """
    if not option_shapes or not options:
        return

    from pptx.util import Inches
    from pptx.enum.text import MSO_ANCHOR

    EMU_PER_IN = 914400
    shapes = sorted(option_shapes, key=lambda s: s.top)[: len(options)]
    texts = [t for t in options][: len(shapes)]
    if not shapes:
        return

    # Trigger: only adapt when an option is long at the DEFAULT template width.
    max_lines_default = max(
        _estimated_wrapped_lines(t, MCQ_OPTION_FONT_PT, MCQ_OPTION_DEFAULT_WIDTH_IN)
        for t in texts
    )
    if max_lines_default < MCQ_OPTION_TRIGGER_LINES:
        return  # short options → keep the exact template layout

    first_top_in = shapes[0].top / EMU_PER_IN
    min_left = min(s.left for s in shapes)
    old_tops = [s.top for s in shapes]

    # Assign each teal A/B/C/D badge (and any per-option pill) to its option by
    # ORIGINAL top BEFORE moving anything. Matching by top-proximity means the
    # tall question text box — which also sits to the LEFT of the options but
    # far ABOVE them — is never mistaken for a badge. Doing the matching up front
    # also avoids double-moving a badge that drifts into a later option's window
    # once options with big deltas slide downward.
    badge_assignment = []  # (badge_shape, option_index)
    tol = Inches(1.2)
    for other in slide.shapes:
        if other in shapes or other.top is None or other.left is None:
            continue
        if other.left >= min_left:
            continue
        nearest = min(range(len(shapes)), key=lambda i: abs(other.top - old_tops[i]))
        if abs(other.top - old_tops[nearest]) <= tol:
            badge_assignment.append((other, nearest))

    # Floor every row to at least the badge height (+ a little) so two badges
    # never pack tighter than their own diameter and visually overlap. Derive
    # this ONLY from the matched badges — NOT every left-side shape — so the
    # (tall, resized) question box can't inflate the floor and push the options
    # clean off the bottom of the slide.
    badge_heights_in = [
        b.height / EMU_PER_IN for b, _ in badge_assignment if b.height is not None
    ]
    min_row_in = max(1.0, (max(badge_heights_in) if badge_heights_in else 0.0) + 0.15)

    def plan(font_pt: int, width_in: float):
        line_h = font_pt * MCQ_OPTION_LINE_HEIGHT_FACTOR / 72.0
        lines = [_estimated_wrapped_lines(t, font_pt, width_in) for t in texts]
        heights = [max(min_row_in, n * line_h) for n in lines]
        tops = [first_top_in]
        for i in range(1, len(shapes)):
            tops.append(tops[i - 1] + heights[i - 1] + MCQ_OPTION_GAP_IN)
        bottom = tops[-1] + heights[-1]
        return tops, heights, bottom

    font_pt = MCQ_OPTION_FONT_PT
    width_in = MCQ_OPTION_WIDENED_WIDTH_IN
    tops, heights, bottom = plan(font_pt, width_in)

    # Last resort only: if the reflow runs off the slide, shrink the font a bit.
    while bottom > MCQ_OPTION_BAND_BOTTOM_IN and font_pt > MCQ_OPTION_MIN_FONT_PT:
        font_pt -= 2
        tops, heights, bottom = plan(font_pt, width_in)

    deltas = []
    for i, sh in enumerate(shapes):
        new_top = int(round(tops[i] * EMU_PER_IN))
        deltas.append(new_top - old_tops[i])
        sh.width = Inches(width_in)
        sh.height = Inches(heights[i])
        sh.text_frame.word_wrap = True
        sh.text_frame.vertical_anchor = MSO_ANCHOR.TOP
        if font_pt != MCQ_OPTION_FONT_PT:
            _set_textbox_font_size(sh, font_pt)
        sh.top = new_top

    for badge, idx in badge_assignment:
        if deltas[idx]:
            badge.top = badge.top + deltas[idx]


def _apply_long_mcq_grid_layout(slide, question_shape, option_shapes, question, options) -> bool:
    """
    Slide-19 type fallback: a *vertical* MCQ whose stem is so long that four
    stacked options would collide with it. When that happens AND all four
    options are short, lay the options out as a 2x2 grid right below the stem:

            A   B
            C   D

    Returns True if the grid layout was applied (caller then skips the normal
    vertical option reflow), False otherwise (caller keeps the vertical layout).

    The teal A/B/C/D badges live as separate shapes to the LEFT of each option
    text box, so each badge is translated by the same delta as its option.
    """
    if not question_shape or not option_shapes or len(option_shapes) < 4:
        return False

    from pptx.util import Inches
    from pptx.enum.text import MSO_ANCHOR

    EMU_PER_IN = 914400

    # ── Trigger 1: stem must be genuinely long ──────────────────────────────
    if len(question) < MCQ_GRID_QUESTION_THRESHOLD:
        return False

    q_font = (
        LONG_MCQ_QUESTION_FONT_PT
        if len(question) < VERY_LONG_MCQ_QUESTION_THRESHOLD
        else VERY_LONG_MCQ_QUESTION_FONT_PT
    )
    q_left_in  = question_shape.left / EMU_PER_IN
    q_top_in   = question_shape.top / EMU_PER_IN
    q_width_in = question_shape.width / EMU_PER_IN
    # Newline-aware: each numbered statement starts on its own line, so a hard
    # break can't be packed onto the previous line. Sum per-segment wraps so we
    # don't undercount (which would let the stem overflow into the first row).
    stem_lines = sum(
        _estimated_wrapped_lines(seg, q_font, q_width_in)
        for seg in (question.split("\n") if "\n" in question else [question])
    )
    if stem_lines < MCQ_GRID_MIN_QUESTION_LINES:
        return False

    # ── Geometry of the two columns ─────────────────────────────────────────
    shapes = sorted(option_shapes, key=lambda s: s.top)[:4]
    texts  = [t for t in options][:4]
    if len(shapes) < 4 or len(texts) < 4:
        return False

    text1_left_in = min(s.left for s in shapes) / EMU_PER_IN
    col_shift_in  = q_width_in / 2.0
    # Column text width: fill half the content band minus the badge gutter on the
    # left and a small right margin, so neither column crosses the centre line.
    text_width_in = max(
        6.0,
        q_width_in / 2.0 - (text1_left_in - q_left_in) - MCQ_GRID_COL_GUTTER_IN,
    )

    # ── Trigger 2: every option must fit the half-width column ──────────────
    if any(
        _estimated_wrapped_lines(t, MCQ_OPTION_FONT_PT, text_width_in)
        > MCQ_GRID_OPTION_MAX_LINES
        for t in texts
    ):
        return False  # long sentence options → grid won't help; keep vertical

    # ── Resize the stem (uncapped vs vertical so options sit truly below) ───
    needed_q_h = min(
        max(stem_lines * q_font * 1.22 / 72.0 + 0.35, 2.6),
        MCQ_GRID_MAX_QUESTION_HEIGHT_IN,
    )
    question_shape.height = Inches(needed_q_h)
    question_shape.text_frame.word_wrap = True
    question_shape.text_frame.vertical_anchor = MSO_ANCHOR.TOP
    _set_textbox_font_size(question_shape, q_font)

    # ── Match each badge to its option BEFORE moving anything ───────────────
    # Badges are the small teal circles sitting just LEFT of each option text
    # box. The (tall) question text box also sits to the left, so we match by
    # top-proximity — its top is far from any option row and is excluded.
    min_left  = min(s.left for s in shapes)
    old_tops  = [s.top for s in shapes]
    old_lefts = [s.left for s in shapes]
    badge_assignment = []  # (badge_shape, option_index)
    tol = Inches(1.2)
    for other in slide.shapes:
        if other in shapes or other.top is None or other.left is None:
            continue
        if other.left >= min_left:
            continue
        nearest = min(range(len(shapes)), key=lambda i: abs(other.top - old_tops[i]))
        if abs(other.top - old_tops[nearest]) <= tol:
            badge_assignment.append((other, nearest))

    # ── Row heights (floored to the badge size so badges never overlap) ─────
    badge_heights_in = [
        b.height / EMU_PER_IN
        for b, _ in badge_assignment
        if b.height is not None
    ]
    min_row_in = max(1.0, (max(badge_heights_in) if badge_heights_in else 0.0) + 0.15)
    line_h = MCQ_OPTION_FONT_PT * MCQ_OPTION_LINE_HEIGHT_FACTOR / 72.0

    def opt_h(text: str) -> float:
        n = _estimated_wrapped_lines(text, MCQ_OPTION_FONT_PT, text_width_in)
        return max(min_row_in, n * line_h)

    # A=idx0, B=idx1, C=idx2, D=idx3  →  grid (col, row)
    grid_pos = {0: (0, 0), 1: (1, 0), 2: (0, 1), 3: (1, 1)}
    row1_h = max(opt_h(texts[0]), opt_h(texts[1]))
    row2_h = max(opt_h(texts[2]), opt_h(texts[3]))
    row1_top_in = q_top_in + needed_q_h + MCQ_GRID_TOP_GAP_IN
    row2_top_in = row1_top_in + row1_h + MCQ_GRID_ROW_GAP_IN

    # ── Place the option text boxes ─────────────────────────────────────────
    deltas = []  # (dleft_emu, dtop_emu) per option
    for i, sh in enumerate(shapes):
        col, row = grid_pos[i]
        new_left_in = text1_left_in + (col_shift_in if col == 1 else 0.0)
        new_top_in  = row1_top_in if row == 0 else row2_top_in
        new_left = int(round(new_left_in * EMU_PER_IN))
        new_top  = int(round(new_top_in * EMU_PER_IN))
        deltas.append((new_left - old_lefts[i], new_top - old_tops[i]))
        sh.left  = new_left
        sh.top   = new_top
        sh.width = Inches(text_width_in)
        sh.height = Inches(row1_h if row == 0 else row2_h)
        sh.text_frame.word_wrap = True
        sh.text_frame.vertical_anchor = MSO_ANCHOR.TOP

    # ── Translate each badge by the same delta as its option ────────────────
    for badge, idx in badge_assignment:
        dleft, dtop = deltas[idx]
        if dleft:
            badge.left = badge.left + dleft
        if dtop:
            badge.top = badge.top + dtop

    return True


def _fill_mcq(slide, content: SlideContent, is_grid: bool = False):
    q = _strip_question_prefix(content.title)
    q = _sanitize_question_title(q)
    question_shapes = _shapes_with_text(slide, "Type question here")
    option_shapes = _shapes_with_text(slide, "Type option here")

    _replace_first(slide, "Type question here", q)
    opts = [_strip_option_prefix(b) for b in content.bullets[:4]]

    if is_grid:
        # In the template the 4 option boxes are arranged COL1-row1, COL1-row2,
        # COL2-row1, COL2-row2 — i.e. A, C, B, D in XML order. We need to fill
        # them in row-major visual order so the user sees A, B, C, D correctly.
        _replace_placeholders_by_shape_position(
            slide, "Type option here", opts, key=_grid_position_key
        )
    else:
        _replace_sequence(slide, "Type option here", opts)
        used_grid = False
        if question_shapes:
            vertical_option_shapes = sorted(option_shapes, key=lambda s: (s.top, s.left))
            # First try the 2-column grid fallback for very long stems with short
            # options (slide-19 type). If it applies, it fully owns the option
            # layout and we skip the vertical reflow below.
            used_grid = _apply_long_mcq_grid_layout(
                slide, question_shapes[0], vertical_option_shapes, q, opts
            )
            if not used_grid:
                _apply_long_vertical_mcq_layout(
                    slide, question_shapes[0], vertical_option_shapes, q
                )
        # After the question band is placed, give long options more width and
        # vertical spacing (only when an option wraps to a 3rd line).
        if not used_grid:
            _apply_long_vertical_mcq_options_layout(slide, option_shapes, opts)

    _clear_unused_placeholders(slide)


def _extract_exam_tag(notes: str) -> str | None:
    """
    Pull ONLY the exam name/year from speaker_notes, stripping any answer
    text that may have leaked onto the same line.

    Expected format:  "Exam: SSC CGL Tier-II 11/09/2019\nAnswer: (a) ..."
    But we also handle:  "Exam: SSC CGL 2019 Answer: (a) Dirge" (no newline).
    """
    if not notes:
        return None
    for line in notes.splitlines():
        line = line.strip()
        if not line.lower().startswith(("exam", "pyq")):
            continue
        # Chop off anything starting with "Answer" on the same line
        import re
        tag = re.split(r'\s*Answer\s*:', line, maxsplit=1, flags=re.IGNORECASE)[0].strip()
        # Remove the "Exam:" prefix itself to keep just the name + year
        for pfx in ("Exam:", "Exam :", "PYQ:", "PYQ :", "Exam-", "Exam"):
            if tag.startswith(pfx):
                tag = tag[len(pfx):].strip()
                break
        if tag:
            # Truncate to keep the banner from overflowing
            if len(tag) > 50:
                tag = tag[:47] + "..."
            return tag
    return None


def _fill_pyq(slide, content: SlideContent, is_grid: bool = False):
    """Same as MCQ but also fills the PYQ subtitle if speaker_notes carries it."""
    _fill_mcq(slide, content, is_grid=is_grid)
    tag = _extract_exam_tag(content.speaker_notes)
    if tag:
        banner_text = f"Question (Exam: {tag})"
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            full = shape.text_frame.text
            if "Question (Type PYQ" in full or "Question (PYQ" in full:
                # set the banner text
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if "Question (Type PYQ" in run.text or "Question (PYQ" in run.text:
                            run.text = banner_text
                # auto-shrink font so the banner never overflows its box
                # default template font is ~32-36pt; reduce for long tags
                total_len = len(banner_text)
                if   total_len > 65: target_pt = 20
                elif total_len > 52: target_pt = 24
                elif total_len > 40: target_pt = 28
                else:                target_pt = None   # keep template default
                if target_pt:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            run.font.size = Pt(target_pt)
                # also enable word-wrap so nothing spills outside the box
                shape.text_frame.word_wrap = True
                break


def _fill_question_only(slide, content: SlideContent):
    from pptx.util import Inches
    from pptx.dml.color import RGBColor

    q = _strip_question_prefix(content.title)
    q = _sanitize_question_title(q)
    _replace_first(slide, "Type question here", q)
    _clear_unused_placeholders(slide)

    # If the writer provided bullets/options (happens when layout was set to
    # question_only but the question actually has MCQ options), render them as
    # a 2-column option grid below the question area so nothing is lost.
    opts = [_strip_option_prefix(b) for b in content.bullets if b.strip()]
    if not opts:
        return

    labels = ["(A)", "(B)", "(C)", "(D)"]
    col_w    = Inches(17.5)
    col_gap  = Inches(3.0)
    left_x   = [Inches(1.5), Inches(1.5) + col_w + col_gap]
    row_h    = Inches(2.8)
    start_y  = Inches(10.5)   # below the question text area

    for i, opt in enumerate(opts[:4]):
        col = i % 2
        row = i // 2
        tb = slide.shapes.add_textbox(
            left_x[col],
            start_y + row * row_h,
            col_w,
            row_h,
        )
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"{labels[i]}  {opt}"
        run.font.size = Pt(34)
        run.font.name = "Poppins"
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def _fill_thank_you(slide, content: SlideContent):
    # The thank-you slide is fully decorative — nothing to replace.
    _clear_unused_placeholders(slide)


def _add_bullets_textbox(slide, bullets, top_in=6.0, font_pt=40,
                         left_in=1.5, width_in=37.0):
    """
    Append a simple bullets textbox below the heading. Used for summary /
    homework where the template only has a small title and no body area.
    All inch coordinates (top_in, left_in, width_in) are given in the 40×22.5
    reference space and are scaled to the current template's canvas here.
    `font_pt` is supplied by the fit engine (reference space); it is also scaled.
    """
    from pptx.util import Inches
    from pptx.dml.color import RGBColor
    sc = _style().scale
    tb = slide.shapes.add_textbox(
        Inches(left_in * sc), Inches(top_in * sc),
        Inches(width_in * sc), Inches(15.0 * sc),
    )
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(max(int(18 * sc), 3))
        run = p.add_run()
        run.text = f"{i + 1}.  {line}"
        run.font.size = Pt(max(int(font_pt * sc), 6))
        run.font.name = "Poppins"
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def _add_summary_illustration(slide, top_after_text: float):
    """
    Place the recap image as a WIDE banner across the BOTTOM of the summary
    slide (about half the slide width), centred horizontally and sitting below
    the last line of text (`top_after_text`) with a margin. Aspect ratio is
    preserved.  All coordinates scale with the current template's canvas.
    """
    from pptx.util import Inches

    if not os.path.exists(_RECAP_ICON_PATH):
        return

    sc = _style().scale
    aspect = 1.834
    try:
        from PIL import Image
        with Image.open(_RECAP_ICON_PATH) as im:
            aspect = im.size[0] / im.size[1]
    except Exception:
        pass

    SLIDE_W   = _style().canvas_w
    SLIDE_BOTTOM = _style().canvas_h - 0.7 * sc
    MARGIN    = 1.0 * sc
    img_w = 18.0 * sc
    img_h = img_w / aspect
    img_l = (SLIDE_W - img_w) / 2.0
    img_t = top_after_text + MARGIN
    img_t = min(img_t, SLIDE_BOTTOM - img_h)

    slide.shapes.add_picture(
        _RECAP_ICON_PATH,
        Inches(img_l), Inches(img_t),
        width=Inches(img_w), height=Inches(img_h),
    )


def _fill_summary_or_homework(slide, content: SlideContent):
    """Template only has heading; add a body textbox for the points."""
    if not content.bullets:
        return
    from pipeline.fit_engine import pick_body_font
    font_pt = pick_body_font(content.bullets, content.layout)

    # Summary slides keep all text full-width at the top, then show the recap
    # image as a wide banner across the bottom below the last line.
    _add_bullets_textbox(slide, content.bullets, top_in=6.0, font_pt=font_pt)

    if content.layout == TemplateType.summary and os.path.exists(_RECAP_ICON_PATH):
        n = len(content.bullets)
        points_bottom = 6.0 + n * (font_pt / 72.0 * 1.25 + 0.25)  # estimated text end
        _add_summary_illustration(slide, points_bottom)


# ─────────────────────────────────────────────────────────────────────────────
# Speaker notes
# ─────────────────────────────────────────────────────────────────────────────

def _set_notes(slide, text):
    if text:
        slide.notes_slide.notes_text_frame.text = text


# ─────────────────────────────────────────────────────────────────────────────
# Router — pick the right filler for each layout
# ─────────────────────────────────────────────────────────────────────────────

def _fill_figure_slide(slide, content: SlideContent):
    """
    Render a detected diagram/figure/formula on a blank dark slide.

    Two modes (set by the teacher during review):
      • image → the cropped PNG, scaled to fit and centred on a white card so a
                white-background crop reads cleanly on the dark canvas.
      • text  → the figure's description rendered as large readable text.

    Everything is a normal movable/resizable shape, so the teacher can reposition
    or delete it in PowerPoint.
    """
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE

    sc       = _style().scale
    canvas_w = _style().canvas_w
    fig = content.figure

    _clear_unused_placeholders(slide)

    # Title pill at the top (reuses the brand yellow-tag style).
    title = content.title or (fig.label if fig else "") or "Diagram"
    try:
        _draw_yellow_title_tag(slide, title, top_in=0.9)
    except Exception:
        pass

    if fig is None:
        return

    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    MUTED = RGBColor(0xC9, 0xCC, 0xD3)

    area_top    = 3.4  * sc
    area_bottom = 20.6 * sc
    image_done = False

    if fig.kind == "image" and fig.image_path and os.path.exists(fig.image_path):
        try:
            from PIL import Image
            with Image.open(fig.image_path) as im:
                iw, ih = im.size
        except Exception:
            iw, ih = 4, 3
        iw = max(1, iw); ih = max(1, ih)

        # Reserve room for a caption line under the image. The user's chosen
        # size scales the max width of the image on its own dedicated slide.
        _OWN_SIZE_MAXW = {"small": 18.0 * sc, "medium": 26.0 * sc, "large": 34.0 * sc}
        max_w = _OWN_SIZE_MAXW.get(getattr(fig, "size", "medium"), 26.0 * sc)
        max_h = (area_bottom - area_top) - 1.6 * sc
        img_scale = min(max_w / iw, max_h / ih)
        w = iw * img_scale
        h = ih * img_scale
        left = (canvas_w - w) / 2.0
        top = area_top

        # White card behind the crop (small padding around the image).
        pad = 0.35 * sc
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left - pad), Inches(top - pad),
            Inches(w + 2 * pad), Inches(h + 2 * pad),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = RGBColor(0xE2, 0xE4, 0xEA)
        card.shadow.inherit = False
        try:
            card.adjustments[0] = 0.04
        except Exception:
            pass

        slide.shapes.add_picture(
            fig.image_path, Inches(left), Inches(top),
            width=Inches(w), height=Inches(h),
        )
        caption_top = top + h + pad + 0.25 * sc
        image_done = True
    else:
        caption_top = area_top

    # Caption / description text.
    cap_parts = []
    if fig.belongs_to:
        cap_parts.append(str(fig.belongs_to))
    if fig.description:
        cap_parts.append(fig.description)
    caption = "  —  ".join(cap_parts) if image_done else (fig.description or "")

    if caption:
        box = slide.shapes.add_textbox(
            Inches(4.0 * sc), Inches(caption_top),
            Inches(canvas_w - 8.0 * sc),
            Inches(max(2.0 * sc, area_bottom - caption_top)),
        )
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP if image_done else MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = caption
        # Smaller for an image caption, larger when text IS the content.
        run.font.size = Pt(max(int(22 * sc), 6) if image_done else max(int(40 * sc), 6))
        run.font.color.rgb = MUTED if image_done else WHITE
        run.font.name = "Poppins"


# Render size → max box dimension (inches) for figures placed ON a slide.
_FIG_SIZE_DIM = {"small": 7.5, "medium": 10.5, "large": 13.5}


def _draw_figure_card(slide, fig, left: float, top: float, max_w: float, max_h: float) -> None:
    """
    Draw ONE figure (image-on-white-card + label, or a text card) fitted inside
    the box (left, top, max_w, max_h) in inches. Everything is a normal movable
    shape so the teacher can reposition/resize it in PowerPoint.
    """
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE

    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    label_h = 0.9 if fig.label else 0.0
    img_zone_h = max_h - label_h

    if fig.kind == "image" and fig.image_path and os.path.exists(fig.image_path):
        try:
            from PIL import Image
            with Image.open(fig.image_path) as im:
                iw, ih = im.size
        except Exception:
            iw, ih = 4, 3
        iw = max(1, iw); ih = max(1, ih)

        pad = 0.25
        avail_w = max_w - 2 * pad
        avail_h = img_zone_h - 2 * pad
        scale = min(avail_w / iw, avail_h / ih)
        w = iw * scale
        h = ih * scale
        card_w = w + 2 * pad
        card_x = left + (max_w - card_w) / 2.0
        card_y = top

        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(card_x), Inches(card_y),
            Inches(card_w), Inches(h + 2 * pad),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = RGBColor(0xE2, 0xE4, 0xEA)
        card.shadow.inherit = False
        try:
            card.adjustments[0] = 0.05
        except Exception:
            pass
        slide.shapes.add_picture(
            fig.image_path, Inches(card_x + pad), Inches(card_y + pad),
            width=Inches(w), height=Inches(h),
        )
        if fig.label:
            chip = slide.shapes.add_textbox(
                Inches(left), Inches(card_y + h + 2 * pad + 0.1),
                Inches(max_w), Inches(label_h),
            )
            tf = chip.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = fig.label
            r.font.size = Pt(16)
            r.font.color.rgb = RGBColor(0xC9, 0xCC, 0xD3)
            r.font.name = "Poppins"
    else:
        # Text figure — a translucent dark card with the description.
        text = fig.description or fig.label or ""
        if not text:
            return
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(top), Inches(max_w), Inches(max_h),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(0x1A, 0x13, 0x10)
        card.line.color.rgb = RGBColor(0xFF, 0xCC, 0x31)
        card.shadow.inherit = False
        tf = card.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.4)
        tf.margin_right = Inches(0.4)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = text
        r.font.size = Pt(20)
        r.font.color.rgb = WHITE
        r.font.name = "Poppins"


def _embed_figures_on_slide(slide, figs) -> None:
    """
    Embed one OR MORE figures directly on a content slide (placement == on_slide).

      • 1 figure  → placed per its `align` (left / center / right) at its `size`.
      • N figures → laid out in a centered horizontal row, each at its `size`,
                    scaled down together if the row would overflow the canvas.

    Every figure is a normal movable shape so the teacher can fine-tune it.
    """
    figs = [f for f in (figs or [])]
    n = len(figs)
    if n == 0:
        return

    sc       = _style().scale
    CANVAS_W = _style().canvas_w
    CANVAS_H = _style().canvas_h

    if n == 1:
        f = figs[0]
        dim = _FIG_SIZE_DIM.get(getattr(f, "size", "medium"), 10.5) * sc
        max_w = dim
        max_h = dim + (0.9 * sc if f.label else 0.0)
        top = 5.0 * sc
        align = getattr(f, "align", "right")
        if align == "left":
            left = 1.0 * sc
        elif align == "center":
            left = (CANVAS_W - max_w) / 2.0
        else:
            left = CANVAS_W - max_w - 1.0 * sc
        _draw_figure_card(slide, f, left, top, max_w, max_h)
        return

    # Multiple figures — centered horizontal row.
    dims = [_FIG_SIZE_DIM.get(getattr(f, "size", "medium"), 10.5) * sc for f in figs]
    gap = 0.8 * sc
    total = sum(dims) + gap * (n - 1)
    avail = CANVAS_W - 2.0 * sc
    if total > avail:
        shrink = avail / total
        dims = [d * shrink for d in dims]
        total = sum(dims) + gap * (n - 1)

    x = (CANVAS_W - total) / 2.0
    row_h = max(dims) + 0.9 * sc
    top = (CANVAS_H - row_h) / 2.0 + 1.5 * sc
    for f, d in zip(figs, dims):
        _draw_figure_card(slide, f, x, top, d, row_h)
        x += d + gap


def _embed_figure_on_slide(slide, fig) -> None:
    """Back-compat single-figure wrapper around _embed_figures_on_slide."""
    _embed_figures_on_slide(slide, [fig])


def _apply_content(slide, content: SlideContent, context: PDFContext, strategy=None):
    t = content.layout
    if t == TemplateType.title_slide:
        _fill_title_slide(slide, content, context)
    elif t == TemplateType.recap_slide:
        _fill_recap_or_topics(slide, content)
    elif t == TemplateType.topics_slide:
        _fill_recap_or_topics(slide, content)
    elif t == TemplateType.section_heading:
        _fill_section_heading(slide, content)
    elif t == TemplateType.theory_slide:
        # theory uses the new yellow-tag + arrow-bullets layout on a blank base
        _fill_theory_slide(slide, content, strategy)
    elif t == TemplateType.table_slide:
        # table-only — yellow caption tag at top, real pptx table fills the body
        _fill_table_slide(slide, content)
    elif t == TemplateType.theory_table_slide:
        # short theory bullets above + table below, non-overlapping
        _fill_theory_table_slide(slide, content, strategy)
    elif t == TemplateType.passage_slide:
        # cloze/comprehension passage — yellow directions banner + verbatim text
        _fill_passage_slide(slide, content)
    elif t == TemplateType.mcq_slide:
        _fill_mcq(slide, content, is_grid=False)
    elif t == TemplateType.mcq_grid_slide:
        _fill_mcq(slide, content, is_grid=True)
    elif t == TemplateType.pyq_slide:
        _fill_pyq(slide, content, is_grid=False)
    elif t == TemplateType.pyq_grid_slide:
        _fill_pyq(slide, content, is_grid=True)
    elif t == TemplateType.question_only:
        _fill_question_only(slide, content)
    elif t == TemplateType.pyq_question_only:
        _fill_question_only(slide, content)
        # patch only the PYQ exam-tag banner — don't re-fill question/options
        tag = _extract_exam_tag(content.speaker_notes)
        if tag:
            banner_text = f"Question (Exam: {tag})"
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                full = shape.text_frame.text
                if "Question (Type PYQ" in full or "Question (PYQ" in full:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if "Question (Type PYQ" in run.text or "Question (PYQ" in run.text:
                                run.text = banner_text
                    total_len = len(banner_text)
                    if   total_len > 65: target_pt = 20
                    elif total_len > 52: target_pt = 24
                    elif total_len > 40: target_pt = 28
                    else:                target_pt = None
                    if target_pt:
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                run.font.size = Pt(target_pt)
                    shape.text_frame.word_wrap = True
                    break
    elif t == TemplateType.summary:
        _fill_first_text(slide, "Summary", content.title or "Summary")
        _fill_summary_or_homework(slide, content)
    elif t == TemplateType.homework_slide:
        _fill_first_text(slide, "Homework", content.title or "Homework")
        _fill_summary_or_homework(slide, content)
    elif t == TemplateType.thank_you_slide:
        _fill_thank_you(slide, content)
    elif t == TemplateType.figure_slide:
        _fill_figure_slide(slide, content)
    else:
        _clear_unused_placeholders(slide)

    _set_notes(slide, content.speaker_notes)


def _fill_first_text(slide, old, new):
    """Replace the first textbox whose text equals `old`."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if shape.text_frame.text.strip() == old:
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    r.text = new
            return True
    return False


# ─────────────────────────────────────────────────────────────────────────────
# MAIN ENTRY
# ─────────────────────────────────────────────────────────────────────────────

def generate_pptx(
    all_slide_contents: list[SlideContent],
    context: PDFContext,
    filename: str = "output.pptx",
    strategy=None,
    template_path: str | None = None,
) -> str:
    """
    Build the final deck by cloning slides from the reference template and
    filling placeholders with our generated content.

    `template_path` overrides the default TEMPLATE_PPTX. All bundled templates
    share the same 14-slot LAYOUT_TO_TEMPLATE_IDX order, so swapping the file
    is the only change needed to switch the visual style.
    """
    os.makedirs(OUTPUT_DIR, exist_ok=True)

    tpl = template_path or TEMPLATE_PPTX

    if not os.path.exists(tpl):
        raise FileNotFoundError(
            f"Reference template not found: {tpl}. "
            "Add the template file to backend/assets/reference_ppts/."
        )

    # Parse the style-guide slide of the chosen template so every _fill_*()
    # function picks up the correct accent colour AND canvas scale for this deck.
    # parse_template_style() caches results, so subsequent calls are instant.
    tpl_basename = os.path.basename(tpl)
    tpl_style = parse_template_style(tpl)
    _tl.style = tpl_style
    print(
        f"  template → {tpl_basename}  "
        f"({tpl_style.canvas_w:.1f}×{tpl_style.canvas_h:.1f}in  "
        f"scale={tpl_style.scale:.3f}  accent=#{tpl_style.accent_hex})"
    )

    prs = Presentation(tpl)
    original_count = len(prs.slides)

    for content in all_slide_contents:
        src_idx = LAYOUT_TO_TEMPLATE_IDX.get(content.layout)
        if src_idx is None:
            # unknown template type — fall back to theory/recap layout
            src_idx = LAYOUT_TO_TEMPLATE_IDX[TemplateType.theory_slide]

        try:
            new_slide = _clone_slide(prs, prs.slides[src_idx])
            _apply_content(new_slide, content, context, strategy)
            # Embed any figures the user pinned ON this slide (placement on_slide),
            # laid out together (1 = aligned/sized, N = a centered sized row).
            _inline = getattr(content, "inline_figures", None) or []
            if _inline:
                try:
                    _embed_figures_on_slide(new_slide, _inline)
                except Exception as fe:
                    print(f"    Slide {content.slide_number:2d} — inline figures failed: {fe}")
            _apply_devanagari_fonts(new_slide)

            print(f"    Slide {content.slide_number:2d} [{content.layout.value:18s}] — "
                  f"{content.title[:55]}")
        except Exception as e:
            print(f"    Slide {content.slide_number:2d} — failed: {e}")

    # Drop the original 14 template slides; keep only the ones we built.
    _delete_slides_by_indices(prs, list(range(original_count)))

    out = os.path.join(OUTPUT_DIR, filename)
    prs.save(out)
    print(f"\n  PPT saved → {out}")
    return out
