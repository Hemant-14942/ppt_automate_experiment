"""
Smoke test: formula rendering in the PPT generator.

Builds a 5-slide mini-deck (no AI, no PDF) with formulas from
physics, chemistry and mathematics to verify:

  1. Unicode conversion  — C6H12O6 → C₆H₁₂O₆, x^2 → x²
  2. Formula renderer    — $$...$$ blocks → matplotlib PNG images
  3. Mixed bullets       — plain text + formulas on same slide
  4. End-to-end output   — valid .pptx written to docs/output/

Run from backend/ with the venv activated:
    python3 _smoke_formula.py

Output file: docs/output/formula_smoke_test.pptx
"""

import sys, os
sys.path.insert(0, os.path.dirname(__file__))

# ── verify formula_renderer standalone ───────────────────────────────────────
print("=" * 60)
print("PHASE 1 — Formula renderer unit checks")
print("=" * 60)

from pipeline.formula_renderer import (
    split_bullet_segments,
    render_formula_png,
    to_unicode_math,
    is_pure_formula_bullet,
    TextSegment,
    FormulaSegment,
)

# 1a. Unicode conversion
tests_unicode = [
    ("H2O",                 "H₂O"),
    ("C6H12O6 + 6O2 -> 6CO2 + 6H2O", "C₆H₁₂O₆ + 6O₂ → 6CO₂ + 6H₂O"),
    ("x^2 + y^2 = r^2",    "x² + y² = r²"),
    ("30 deg C",             "30° C"),
    ("E = mc^2",             "E = mc²"),
]
all_ok = True
for raw, expected in tests_unicode:
    got = to_unicode_math(raw)
    ok = got == expected
    if not ok:
        all_ok = False
    status = "✓" if ok else "✗"
    print(f"  {status}  {raw!r}  →  {got!r}  {'(expected: ' + repr(expected) + ')' if not ok else ''}")

# 1b. Bullet segmentation
print()
print("Segmentation tests:")
seg_tests = [
    ("$$E = mc^2$$",               [(FormulaSegment, True)]),
    ("Speed: $c = 3\\times10^8$ m/s", [(TextSegment,), (FormulaSegment, False), (TextSegment,)]),
    ("No formula here",             [(TextSegment,)]),
]
for text, expected_types in seg_tests:
    segs = split_bullet_segments(text)
    types = [(type(s).__name__) for s in segs]
    print(f"  {text!r}")
    print(f"    → {types}")

# 1c. Render a sample formula
print()
print("Rendering $$\\frac{-b \\pm \\sqrt{b^2 - 4ac}}{2a}$$ …")
png = render_formula_png(r"\frac{-b \pm \sqrt{b^2 - 4ac}}{2a}", fontsize=32)
if png:
    print(f"  ✓  Got PNG ({len(png):,} bytes)")
else:
    print("  ✗  Render returned None — check matplotlib install")
    sys.exit(1)

# ── Build PPT ────────────────────────────────────────────────────────────────
print()
print("=" * 60)
print("PHASE 2 — Building smoke-test PPTX")
print("=" * 60)

from schemas.slide_content import SlideContent
from schemas.slide_plan import TemplateType

# Slide definitions — each entry is kwargs to SlideContent
SLIDES = [
    # --- Slide 1: title ---
    dict(
        slide_number=1,
        title="Formula Rendering — Smoke Test",
        bullets=[
            "Auto-converts chemistry subscripts: H2SO4, C6H6",
            "Renders math blocks as images: $$E = mc^2$$",
            "Unicode: x^2 + y^2 = r^2, 30 deg C",
        ],
        speaker_notes="Smoke test title slide",
        layout=TemplateType.theory_slide,
    ),

    # --- Slide 2: physics ---
    dict(
        slide_number=2,
        title="Physics — Resonance Frequency",
        bullets=[
            "Resonant frequency of an LC circuit:",
            "$$f_0 = \\frac{1}{2\\pi\\sqrt{LC}}$$",
            "where L = inductance (Henry), C = capacitance (Farad)",
            "At resonance: X_L = X_C, impedance is purely resistive",
            "Power factor: $$\\cos\\phi = \\frac{R}{Z}$$",
        ],
        speaker_notes="LC circuit resonance frequency, physics.",
        layout=TemplateType.theory_slide,
    ),

    # --- Slide 3: chemistry ---
    dict(
        slide_number=3,
        title="Chemistry — Combustion of Glucose",
        bullets=[
            "Complete combustion reaction:",
            "C6H12O6 + 6O2 -> 6CO2 + 6H2O  (exothermic)",
            "Enthalpy change (standard conditions):",
            "$$\\Delta H = -2803 \\text{ kJ/mol}$$",
            "Avogadro's number: N_A = 6.022 x 10^23 mol^-1",
            "Molar mass of glucose: 180.16 g/mol",
        ],
        speaker_notes="Glucose combustion chemistry.",
        layout=TemplateType.theory_slide,
    ),

    # --- Slide 4: maths ---
    dict(
        slide_number=4,
        title="Mathematics — Quadratic Formula",
        bullets=[
            "For ax^2 + bx + c = 0, the roots are:",
            "$$x = \\frac{-b \\pm \\sqrt{b^2 - 4ac}}{2a}$$",
            "Discriminant D = b^2 - 4ac determines nature of roots:",
            "(a) D > 0  →  two distinct real roots",
            "(b) D = 0  →  one repeated real root",
            "(c) D < 0  →  two complex conjugate roots",
            "Sum of roots: $$\\alpha + \\beta = \\frac{-b}{a}$$",
            "Product of roots: $$\\alpha \\cdot \\beta = \\frac{c}{a}$$",
        ],
        speaker_notes="Quadratic formula and discriminant.",
        layout=TemplateType.theory_slide,
    ),

    # --- Slide 5: mixed inline + block ---
    dict(
        slide_number=5,
        title="Mixed Formulas — Physics & Chemistry",
        bullets=[
            "Newton's law: F = ma (where m in kg, a in m/s^2)",
            "$$F = G\\frac{m_1 m_2}{r^2}$$  (Universal gravitation)",
            "Ideal gas: PV = nRT  →  $$P = \\frac{nRT}{V}$$",
            "Chemical equilibrium: H2 + I2 <-> 2HI",
            "pH = -log[H+], for HCl: [H+] = 0.1 M → pH ≈ 1",
            "$$\\lambda = \\frac{h}{mv}$$  (de Broglie wavelength)",
        ],
        speaker_notes="Mixed physics and chemistry formulas.",
        layout=TemplateType.theory_slide,
    ),
]

# Build SlideContent objects — the sanitizer will apply Unicode conversion
slides = [SlideContent(**s) for s in SLIDES]

print("Slide content (after Unicode auto-conversion):")
for sc in slides:
    print(f"\n  Slide {sc.slide_number}: {sc.title}")
    for b in sc.bullets:
        print(f"    • {b}")

# ── Generate PPTX ─────────────────────────────────────────────────────────────
from config import TEMPLATE_PPTX, OUTPUT_DIR
from pptx import Presentation
from schemas.request import PDFContext

# Use the first available reference template
import glob as _glob
tpl_files = _glob.glob(os.path.join(os.path.dirname(TEMPLATE_PPTX), "*.pptx"))
if not tpl_files:
    print("\n✗  No template .pptx found — cannot build PPT")
    sys.exit(1)
tpl_path = tpl_files[0]
print(f"\nUsing template: {os.path.basename(tpl_path)}")

from pipeline.ppt_generator import _apply_content, _clone_slide, LAYOUT_TO_TEMPLATE_IDX
from pptx.util import Inches, Pt

prs = Presentation(tpl_path)
ref_prs = Presentation(tpl_path)   # read-only reference for cloning

context = PDFContext(
    batch="Smoke Test Batch",
    purpose="Lecture notes",
    subject="Science",
    class_level="Class 11-12",
)

for sc in slides:
    layout_idx = LAYOUT_TO_TEMPLATE_IDX.get(sc.layout, 3)   # blank dark slide
    src_slide  = ref_prs.slides[layout_idx]
    new_slide  = _clone_slide(prs, src_slide)
    _apply_content(new_slide, sc, context)
    print(f"  Slide {sc.slide_number} — {sc.layout.value}  ✓")

# Remove template placeholder slides (the reference deck may have them)
# Keep only our newly added slides (last N)
n_orig = len(ref_prs.slides)
n_total = len(prs.slides)
n_new = n_total - n_orig

if n_new > 0:
    from pipeline.ppt_generator import _delete_slides_by_indices
    _delete_slides_by_indices(prs, list(range(n_orig)))

os.makedirs(OUTPUT_DIR, exist_ok=True)
out_path = os.path.join(OUTPUT_DIR, "formula_smoke_test.pptx")
prs.save(out_path)
print(f"\n✓  Saved to: {out_path}")
print(f"   Size: {os.path.getsize(out_path):,} bytes")
print()
print("Open the file in PowerPoint / LibreOffice to verify formulas render as images.")
