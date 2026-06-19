"""
One-time migration script: restructure Clat & Architecture templates
to match Common Template's 14-slide slot order.

Run from backend/:
    python migrate_templates.py

SLOT ORDER (matching Common Template exactly):
  0  recap_slide
  1  topics_slide
  2  section_heading
  3  theory_slide   (blank dark)
  4  mcq_slide      (vertical A/B/C/D)
  5  mcq_grid_slide (2x2 grid)
  6  question_only
  7  pyq_slide      (vertical + PYQ bar)
  8  pyq_grid_slide (2x2 + PYQ bar)
  9  pyq_question_only
 10  summary
 11  homework_slide
 12  thank_you_slide
 13  style_guide    (skip slide)
"""
import os
import copy
import shutil
from lxml import etree
from pptx import Presentation
from pptx.util import Pt, Emu, Inches
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

REFS = os.path.join(os.path.dirname(__file__), "assets", "reference_ppts")


# ─────────────────────────────────────────────────────────────────────────────
# Low-level clone helper (same logic as ppt_generator._clone_slide)
# ─────────────────────────────────────────────────────────────────────────────

def _clone_slide(prs, src_slide):
    """Deep-copy a slide (shapes + background + media) into a new blank slide."""
    blank_layout = prs.slide_layouts[0]
    new_slide = prs.slides.add_slide(blank_layout)

    new_cSld = new_slide._element.find(qn("p:cSld"))
    src_cSld = src_slide._element.find(qn("p:cSld"))
    new_sptree = new_cSld.find(qn("p:spTree"))
    src_sptree = src_cSld.find(qn("p:spTree"))

    # Copy shapes
    for child in list(new_sptree):
        if etree.QName(child).localname not in ("nvGrpSpPr", "grpSpPr"):
            new_sptree.remove(child)
    for child in src_sptree:
        if etree.QName(child).localname in ("nvGrpSpPr", "grpSpPr"):
            continue
        new_sptree.append(copy.deepcopy(child))

    # Copy background
    src_bg = src_cSld.find(qn("p:bg"))
    if src_bg is not None:
        existing_bg = new_cSld.find(qn("p:bg"))
        if existing_bg is not None:
            new_cSld.remove(existing_bg)
        new_cSld.insert(list(new_cSld).index(new_sptree), copy.deepcopy(src_bg))

    # Remap relationship IDs for embedded images / media
    R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    src_part = src_slide.part
    new_part = new_slide.part
    rid_map: dict[str, str] = {}

    def _ensure_rid(old_rid):
        if old_rid in rid_map:
            return rid_map[old_rid]
        try:
            rel = src_part.rels[old_rid]
        except KeyError:
            return None
        new_rid = new_part.relate_to(rel.target_part, rel.reltype)
        rid_map[old_rid] = new_rid
        return new_rid

    for elem in new_cSld.iter():
        for attr in (f"{{{R_NS}}}embed", f"{{{R_NS}}}link"):
            old = elem.get(attr)
            if old:
                new = _ensure_rid(old)
                if new:
                    elem.set(attr, new)

    return new_slide


def _delete_slides_by_indices(prs, indices):
    """Drop slides at given 0-based indices."""
    xml_slides = prs.slides._sldIdLst
    slides_list = list(xml_slides)
    for idx in sorted(indices, reverse=True):
        if 0 <= idx < len(slides_list):
            sl_el = slides_list[idx]
            rid = sl_el.get(qn("r:id"))
            xml_slides.remove(sl_el)
            prs.part.drop_rel(rid)


def _reorder_slides(prs, new_order):
    """
    Reorder slides in-place. new_order is a list of CURRENT 0-based indices
    representing the desired final order.
    e.g. [2, 0, 1] → current slide 2 first, then 0, then 1.
    """
    xml_slides = prs.slides._sldIdLst
    slide_elements = list(xml_slides)
    # remove all
    for el in slide_elements:
        xml_slides.remove(el)
    # re-append in desired order
    for idx in new_order:
        xml_slides.append(slide_elements[idx])


# ─────────────────────────────────────────────────────────────────────────────
# Text helpers
# ─────────────────────────────────────────────────────────────────────────────

def _first_text_shape(slide):
    """Return the first shape that has a text frame."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            return shape
    return None


def _set_run_text(run, text, font_name=None, size_pt=None, bold=None, rgb=None):
    run.text = text
    if font_name:
        run.font.name = font_name
    if size_pt:
        run.font.size = Pt(size_pt)
    if bold is not None:
        run.font.bold = bold
    if rgb:
        run.font.color.rgb = RGBColor(*bytes.fromhex(rgb.lstrip("#")))


def _add_textbox(slide, left, top, width, height,
                 text, font_name, size_pt, bold, rgb_hex):
    """Add a new text box to a slide with given position and style."""
    from pptx.util import Pt, Emu
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    run = para.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*bytes.fromhex(rgb_hex.lstrip("#")))
    return txBox


# ─────────────────────────────────────────────────────────────────────────────
# Per-template PYQ bar addition
# ─────────────────────────────────────────────────────────────────────────────

def _add_pyq_bar(slide, accent_hex, font_name="Poppins", size_pt=13):
    """
    Add a PYQ exam/year info text box at the top of the slide.
    Positioned to sit above the question text area.
    """
    # Slide canvas is 10 inches wide, 5.625 inches tall (widescreen 16:9 standard)
    left   = Inches(0.3)
    top    = Inches(0.18)
    width  = Inches(9.4)
    height = Inches(0.45)
    _add_textbox(
        slide, left, top, width, height,
        text="(Type PYQ Exam & Year Info For PYQ Questions)",
        font_name=font_name,
        size_pt=size_pt,
        bold=True,
        rgb_hex=accent_hex,
    )


# ─────────────────────────────────────────────────────────────────────────────
# CLAT EVENING FORMAT MIGRATION
# ─────────────────────────────────────────────────────────────────────────────

def migrate_clat(path):
    print(f"\n{'='*55}")
    print(f"  Migrating: {os.path.basename(path)}")
    print(f"{'='*55}")

    prs = Presentation(path)
    n = len(prs.slides)
    print(f"  Original slide count: {n}")

    # Clat original slots:
    #  [0] Style guide
    #  [1] Section heading
    #  [2] MCQ vertical (A/B/C/D stacked)
    #  [3] MCQ 2x2 grid
    #  [4] Question only
    #  [5] Blank content slide

    # ── Step 1: Clone new slides (appended at end, will be reordered) ────────

    # [6 new] Recap — clone section heading [1], change text
    recap = _clone_slide(prs, prs.slides[1])
    for shape in recap.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if "Type Heading Here" in run.text:
                        run.text = "Recap"
    # Add subtitle
    _add_textbox(recap,
        left=Inches(0.5), top=Inches(1.6), width=Inches(9), height=Inches(0.7),
        text="of previous lecture",
        font_name="Poppins", size_pt=22, bold=False, rgb_hex="#FFFFFF")
    print("  Created: Recap slide")

    # [7 new] Topics — same base
    topics = _clone_slide(prs, prs.slides[1])
    for shape in topics.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if "Type Heading Here" in run.text:
                        run.text = "Topics"
    _add_textbox(topics,
        left=Inches(0.5), top=Inches(1.6), width=Inches(9), height=Inches(0.7),
        text="to be covered",
        font_name="Poppins", size_pt=22, bold=False, rgb_hex="#FFFFFF")
    print("  Created: Topics slide")

    # [8 new] PYQ MCQ vertical — clone MCQ vertical [2], add PYQ bar
    pyq_vert = _clone_slide(prs, prs.slides[2])
    _add_pyq_bar(pyq_vert, accent_hex="#FFCC31", font_name="Poppins")
    print("  Created: PYQ MCQ vertical")

    # [9 new] PYQ 2x2 grid — clone MCQ grid [3], add PYQ bar
    pyq_grid = _clone_slide(prs, prs.slides[3])
    _add_pyq_bar(pyq_grid, accent_hex="#FFCC31", font_name="Poppins")
    print("  Created: PYQ 2x2 grid")

    # [10 new] PYQ question only — clone question only [4], add PYQ bar
    pyq_qonly = _clone_slide(prs, prs.slides[4])
    _add_pyq_bar(pyq_qonly, accent_hex="#FFCC31", font_name="Poppins")
    print("  Created: PYQ question only")

    # [11 new] Summary — clone section heading [1]
    summary = _clone_slide(prs, prs.slides[1])
    for shape in summary.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if "Type Heading Here" in run.text:
                        run.text = "Summary"
    print("  Created: Summary")

    # [12 new] Homework — clone section heading [1]
    homework = _clone_slide(prs, prs.slides[1])
    for shape in homework.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if "Type Heading Here" in run.text:
                        run.text = "Homework"
    print("  Created: Homework")

    # [13 new] Thank You — clone blank slide [5], add text
    thankyou = _clone_slide(prs, prs.slides[5])
    _add_textbox(thankyou,
        left=Inches(2.5), top=Inches(2.0), width=Inches(5), height=Inches(1.5),
        text="Thank You",
        font_name="Poppins ExtraBold", size_pt=48, bold=True, rgb_hex="#FFFFFF")
    print("  Created: Thank You")

    # At this point slide indices are:
    #   0=StyleGuide  1=SectionHeading  2=MCQvert  3=MCQgrid  4=QOnly  5=Blank
    #   6=Recap  7=Topics  8=PYQvert  9=PYQgrid  10=PYQqonly  11=Summary
    #   12=Homework  13=ThankYou
    #
    # Target order (14 slots):
    #   0=Recap        → idx 6
    #   1=Topics       → idx 7
    #   2=SectionHdg   → idx 1
    #   3=Blank        → idx 5
    #   4=MCQvert      → idx 2
    #   5=MCQgrid      → idx 3
    #   6=QOnly        → idx 4
    #   7=PYQvert      → idx 8
    #   8=PYQgrid      → idx 9
    #   9=PYQqonly     → idx 10
    #  10=Summary      → idx 11
    #  11=Homework     → idx 12
    #  12=ThankYou     → idx 13
    #  13=StyleGuide   → idx 0

    new_order = [6, 7, 1, 5, 2, 3, 4, 8, 9, 10, 11, 12, 13, 0]
    _reorder_slides(prs, new_order)
    print(f"  Reordered → {len(prs.slides)} slides")

    prs.save(path)
    print(f"  Saved → {path}")


# ─────────────────────────────────────────────────────────────────────────────
# ARCHITECTURE FORMAT MIGRATION
# ─────────────────────────────────────────────────────────────────────────────

def migrate_architecture(path):
    print(f"\n{'='*55}")
    print(f"  Migrating: {os.path.basename(path)}")
    print(f"{'='*55}")

    prs = Presentation(path)
    n = len(prs.slides)
    print(f"  Original slide count: {n}")

    # Architecture original slots:
    #  [0] Recap
    #  [1] Topics
    #  [2] Style guide
    #  [3] Section heading
    #  [4] Blank content
    #  [5] MCQ vertical
    #  [6] MCQ 2x2 grid
    #  [7] Question only
    #  [8] Thank You

    # ── Step 1: Clone new slides ─────────────────────────────────────────────

    # [9 new] PYQ MCQ vertical — clone MCQ vertical [5], add PYQ bar
    pyq_vert = _clone_slide(prs, prs.slides[5])
    _add_pyq_bar(pyq_vert, accent_hex="#FFC000", font_name="Poppins")
    print("  Created: PYQ MCQ vertical")

    # [10 new] PYQ 2x2 grid — clone MCQ grid [6], add PYQ bar
    pyq_grid = _clone_slide(prs, prs.slides[6])
    _add_pyq_bar(pyq_grid, accent_hex="#FFC000", font_name="Poppins")
    print("  Created: PYQ 2x2 grid")

    # [11 new] PYQ question only — clone question only [7], add PYQ bar
    pyq_qonly = _clone_slide(prs, prs.slides[7])
    _add_pyq_bar(pyq_qonly, accent_hex="#FFC000", font_name="Poppins")
    print("  Created: PYQ question only")

    # [12 new] Summary — clone section heading [3]
    summary = _clone_slide(prs, prs.slides[3])
    for shape in summary.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if "Type Heading Here" in run.text:
                        run.text = "Summary"
    print("  Created: Summary")

    # [13 new] Homework — clone section heading [3]
    homework = _clone_slide(prs, prs.slides[3])
    for shape in homework.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if "Type Heading Here" in run.text:
                        run.text = "Homework"
    print("  Created: Homework")

    # At this point slide indices are:
    #   0=Recap  1=Topics  2=StyleGuide  3=SectionHdg  4=Blank
    #   5=MCQvert  6=MCQgrid  7=QOnly  8=ThankYou
    #   9=PYQvert  10=PYQgrid  11=PYQqonly  12=Summary  13=Homework
    #
    # Target order:
    #   0=Recap        → idx 0
    #   1=Topics       → idx 1
    #   2=SectionHdg   → idx 3
    #   3=Blank        → idx 4
    #   4=MCQvert      → idx 5
    #   5=MCQgrid      → idx 6
    #   6=QOnly        → idx 7
    #   7=PYQvert      → idx 9
    #   8=PYQgrid      → idx 10
    #   9=PYQqonly     → idx 11
    #  10=Summary      → idx 12
    #  11=Homework     → idx 13
    #  12=ThankYou     → idx 8
    #  13=StyleGuide   → idx 2

    new_order = [0, 1, 3, 4, 5, 6, 7, 9, 10, 11, 12, 13, 8, 2]
    _reorder_slides(prs, new_order)
    print(f"  Reordered → {len(prs.slides)} slides")

    prs.save(path)
    print(f"  Saved → {path}")


# ─────────────────────────────────────────────────────────────────────────────
# Verification
# ─────────────────────────────────────────────────────────────────────────────

def verify(path):
    prs = Presentation(path)
    print(f"\n  Verification: {os.path.basename(path)} ({len(prs.slides)} slides)")
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    texts.append(t[:50])
        label = " | ".join(texts[:2]) if texts else "(no text)"
        print(f"    [{i:2d}]  {label}")


# ─────────────────────────────────────────────────────────────────────────────
# Main
# ─────────────────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    clat_path = os.path.join(REFS, "Clat evening format.pptx")
    arch_path = os.path.join(REFS, "Acchitecture Format.pptx")

    # Backup originals
    for p in [clat_path, arch_path]:
        bak = p + ".bak"
        if not os.path.exists(bak):
            shutil.copy2(p, bak)
            print(f"  Backed up → {os.path.basename(bak)}")

    migrate_clat(clat_path)
    migrate_architecture(arch_path)

    print("\n\n── Final verification ──────────────────────────────────────")
    verify(clat_path)
    verify(arch_path)

    print("\n\nDone! Both templates now have 14 slides in Common Template order.\n")
